"""Fix slide issues: remove duplicate Galileo slides and recreate missing Einstein slides."""

from pptx import Presentation
from copy import deepcopy
from lxml import etree
import warnings

warnings.filterwarnings("ignore")

NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
NSMAP = {"a": NS_A}

# Colors
WHITE = "E0E0E0"
CYAN = "5CC9F5"
GRAY = "808888"
GOLD = "FFD54F"
RED = "FF8A80"
GREEN = "69F0AE"
ORANGE = "FFAB40"
HEADER = "FFFFFF"


def colorize_line(line):
    stripped = line.strip()
    if not stripped:
        return [("", WHITE, False)]
    if "═══" in stripped:
        return [(line, HEADER, True)]
    if stripped.startswith("#") and "═══" not in stripped:
        if "CONTRADICTION" in stripped or "矛盾" in stripped:
            return [(line, ORANGE, False)]
        if "★" in stripped:
            return [(line, GOLD, True)]
        if "⚠" in stripped:
            return [(line, ORANGE, False)]
        return [(line, GRAY, False)]
    if stripped.startswith("$ gaia"):
        parts = []
        if "  #" in line:
            idx = line.index("  #")
            parts.append((line[:idx], CYAN, False))
            parts.append((line[idx:], GRAY, False))
        elif "# " in line and not line.strip().startswith("#"):
            idx = line.index("# ")
            parts.append((line[:idx], CYAN, False))
            parts.append((line[idx:], GRAY, False))
        else:
            parts.append((line, CYAN, False))
        return parts
    if "★" in line:
        return [(line, GOLD, True)]
    if "contradiction" in stripped.lower() or "CONTRADICTION" in stripped:
        return [(line, ORANGE, True)]
    if stripped.startswith("→") or stripped.startswith("→"):
        if "★" in line:
            return [(line, GOLD, True)]
        return [(line, GREEN, False)]
    has_down = "↓" in line
    has_up = "↑" in line
    if has_down and has_up:
        segments = []
        parts = line.split("  ")
        for part in parts:
            ps = part.strip()
            if not ps:
                continue
            if "↓" in part:
                segments.append(("  " + ps, RED, False))
            elif "↑" in part:
                segments.append(("  " + ps, GREEN, False))
            else:
                segments.append(("  " + ps, WHITE, False))
        return segments if segments else [(line, WHITE, False)]
    if has_down:
        return [(line, RED, False)]
    if has_up:
        return [(line, GREEN, False)]
    return [(line, WHITE, False)]


def make_run_xml(text, color_hex, bold=False):
    ns = NS_A
    r = etree.Element(f"{{{ns}}}r")
    rPr = etree.SubElement(r, f"{{{ns}}}rPr")
    rPr.set("lang", "en-US")
    rPr.set("dirty", "0")
    if bold:
        rPr.set("b", "1")
    solidFill = etree.SubElement(rPr, f"{{{ns}}}solidFill")
    srgbClr = etree.SubElement(solidFill, f"{{{ns}}}srgbClr")
    srgbClr.set("val", color_hex)
    latin = etree.SubElement(rPr, f"{{{ns}}}latin")
    latin.set("typeface", "Helvetica Neue")
    t = etree.SubElement(r, f"{{{ns}}}t")
    t.text = text
    if text and (text[0] == " " or text[-1] == " "):
        t.set("{http://www.w3.org/XML/1998/namespace}space", "preserve")
    return r


def set_colored_text(text_frame, lines):
    ref_p = text_frame.paragraphs[0]._p
    ref_pPr = ref_p.find("a:pPr", NSMAP)
    ref_pPr_xml = deepcopy(ref_pPr) if ref_pPr is not None else None

    txBody = text_frame._txBody
    for p in txBody.findall(f"{{{NS_A}}}p"):
        txBody.remove(p)

    for i, line in enumerate(lines):
        p = etree.SubElement(txBody, f"{{{NS_A}}}p")
        if ref_pPr_xml is not None:
            pPr = deepcopy(ref_pPr_xml)
            if i > 0 and "algn" in pPr.attrib:
                del pPr.attrib["algn"]
            p.insert(0, pPr)
        segments = colorize_line(line)
        for text, color, bold in segments:
            r = make_run_xml(text, color, bold)
            p.append(r)


def duplicate_slide(prs, source_idx):
    """Duplicate a slide to the end."""
    source = prs.slides[source_idx]
    new_slide = prs.slides.add_slide(source.slide_layout)
    sp_tree = new_slide.shapes._spTree
    for child in list(sp_tree)[2:]:
        sp_tree.remove(child)
    src_tree = source.shapes._spTree
    for child in list(src_tree)[2:]:
        sp_tree.append(deepcopy(child))
    return new_slide


def move_slide(prs, from_idx, to_idx):
    sldIdLst = prs.slides._sldIdLst
    slides = list(sldIdLst)
    el = slides[from_idx]
    sldIdLst.remove(el)
    slides = list(sldIdLst)
    if to_idx >= len(slides):
        sldIdLst.append(el)
    else:
        slides[to_idx].addprevious(el)


def delete_slide(prs, index):
    """Properly delete a slide by index."""
    # Remove from slide id list
    sldIdLst = prs.slides._sldIdLst
    slides = list(sldIdLst)
    el = slides[index]
    sldIdLst.remove(el)

    # Also remove the relationship
    rId = el.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    if rId:
        try:
            prs.part.drop_rel(rId)
        except Exception:
            pass  # Some versions don't support this


# Einstein slide content
EINSTEIN_OVERVIEW_TEXT = [
    "1919 年日食远征: 观测证实广义相对论",
    "",
    "5 个 knowledge package，跨越 200 年:",
    "  Pkg 1  先验知识 (Newton, Maxwell, Eötvös)",
    "  Pkg 2  等效原理 (电梯思想实验)",
    "  Pkg 3  1911 光偏折预测 (0.87″)",
    "  Pkg 4  广义相对论 → 1.75″ (定量矛盾!)",
    "  Pkg 5  Eddington 日食观测 → 判决",
    "",
    "展示 Gaia 如何处理:",
    "  • 同一现象的竞争理论 (Newton vs GR)",
    "  • 定量矛盾 (0.87″ vs 1.75″)",
    "  • 观测证据的判决性作用",
]

EINSTEIN_PKG13_CODE = [
    "$ gaia init einstein_elevator",
    "",
    "# ═══ Pkg 1: prior_knowledge ═══",
    '$ gaia node add "F=mᵢa" --prior 0.95        → 6001',
    '$ gaia node add "F=GMm/r²" --prior 0.95     → 6002',
    '$ gaia node add "光的微粒说" --prior 0.50     → 6003',
    '$ gaia node add "Maxwell EM" --prior 0.90    → 6005',
    '$ gaia node add "Eötvös mᵢ=mᵍ" --prior 0.95 → 6006',
    "$ gaia edge add --tail 6001,6002,6003 \\",
    '    --head 6004 --type deduction  P=0.60',
    '  → 6004 "Soldner: 0.87″" prior=0.60',
    "",
    "# ═══ Pkg 2: equivalence_principle ═══",
    '$ gaia node add "电梯思想实验" --prior 1.0  → 6007',
    "$ gaia edge add --tail 6006,6007 --head 6008",
    '  → 6008 "等效原理" belief≈0.85',
    "$ gaia edge add --tail 6008,6005 --head 6009",
    '  → 6009 "光在引力场弯曲"',
    "",
    "# ═══ Pkg 3: 1911_light_deflection ═══",
    "$ gaia edge add --tail 6009 --head 6010",
    '  → 6010 "Einstein: 0.87″" belief≈0.80',
    "$ gaia commit && gaia propagate",
    "# ⚠ 6004 = 6010 = 0.87″  此时无矛盾!",
]


def main():
    prs = Presentation("Gaia_Presentation.pptx")
    total = len(prs.slides)
    print(f"Starting with {total} slides")

    # Step 1: Remove duplicate Galileo slides at idx 15 and 16
    # Delete from higher index first
    delete_slide(prs, 16)
    delete_slide(prs, 15)
    print(f"Removed duplicates. Now: {len(prs.slides)} slides")

    # Step 2: Create 2 new Einstein slides (overview + Pkg 1-3)
    # Duplicate from the existing Einstein Pkg 4 slide (now at idx 15)
    # to get the right layout/template
    einstein_pkg4_idx = 15  # After removing 2 slides, Einstein Pkg 4 is at idx 15

    duplicate_slide(prs, einstein_pkg4_idx)  # Einstein overview
    duplicate_slide(prs, einstein_pkg4_idx)  # Einstein Pkg 1-3

    # Move them to position 15 (before Pkg 4)
    total = len(prs.slides)
    move_slide(prs, total - 1, 15)  # Move Pkg 1-3 to idx 15
    total_after = len(prs.slides)
    move_slide(prs, total_after - 1, 15)  # Move Overview to idx 15

    print(f"Added Einstein slides. Now: {len(prs.slides)} slides")

    # Step 3: Update the 2 new Einstein slides
    # Slide at idx 15: Einstein Overview
    slide_ein_overview = prs.slides[15]
    # Update Part header
    for shape in slide_ein_overview.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text
            if "Part 3" in text:
                shape.text_frame.paragraphs[0].runs[0].text = "Part 4 · 爱因斯坦思想实验"
            break
    # Update subtitle
    slide_ein_overview.shapes[1].text_frame.paragraphs[0].runs[0].text = "爱因斯坦电梯 — 概述"
    # Update code block with overview text
    set_colored_text(slide_ein_overview.shapes[3].text_frame, EINSTEIN_OVERVIEW_TEXT)
    # Clear right-side shapes (they're from the template, not relevant)
    clear_right_shapes(slide_ein_overview)

    # Slide at idx 16: Einstein Pkg 1-3
    slide_ein_pkg13 = prs.slides[16]
    for shape in slide_ein_pkg13.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text
            if "Part 3" in text:
                shape.text_frame.paragraphs[0].runs[0].text = "Part 4 · 爱因斯坦思想实验"
            break
    slide_ein_pkg13.shapes[1].text_frame.paragraphs[0].runs[0].text = (
        "Pkg 1-3: 先验知识 → 等效原理 → 1911 预测"
    )
    set_colored_text(slide_ein_pkg13.shapes[3].text_frame, EINSTEIN_PKG13_CODE)
    clear_right_shapes(slide_ein_pkg13)

    # Step 4: Verify and save
    prs.save("Gaia_Presentation.pptx")

    # Verify
    prs2 = Presentation("Gaia_Presentation.pptx")
    print(f"\nFinal slide count: {len(prs2.slides)}")
    for i in range(len(prs2.slides)):
        slide = prs2.slides[i]
        header = ''
        subtitle = ''
        for j, shape in enumerate(slide.shapes):
            if shape.has_text_frame:
                if j == 0:
                    header = shape.text_frame.text[:30]
                if j == 1:
                    subtitle = shape.text_frame.text[:50]
        print(f"  {i+1}. [{header}] {subtitle}")


def clear_right_shapes(slide, x_threshold=5400000):
    sp_tree = slide.shapes._spTree
    to_remove = []
    for child in list(sp_tree)[2:]:
        off = child.find(f".//{{{NS_A}}}off")
        if off is not None:
            x = int(off.get("x", "0"))
            if x >= x_threshold:
                to_remove.append(child)
    for child in to_remove:
        sp_tree.remove(child)


if __name__ == "__main__":
    main()
