"""Rewrite Galileo slides (11-13) with proper gaia CLI commands."""

from pptx import Presentation
from copy import deepcopy
from lxml import etree


def clear_and_set_text(text_frame, lines):
    """Replace all text in a text_frame with new lines."""
    ref_para = text_frame.paragraphs[0]._p
    ref_xml = deepcopy(ref_para)
    nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}

    txBody = text_frame._txBody
    for p in txBody.findall(f'{{{nsmap["a"]}}}p'):
        txBody.remove(p)

    for i, line in enumerate(lines):
        new_p = deepcopy(ref_xml)
        for r in new_p.findall("a:r", nsmap):
            new_p.remove(r)
        pPr = new_p.find("a:pPr", nsmap)
        if pPr is not None and i > 0:
            if "algn" in pPr.attrib:
                del pPr.attrib["algn"]

        r_elem = deepcopy(ref_xml.findall("a:r", nsmap)[0]) if ref_xml.findall("a:r", nsmap) else None
        if r_elem is not None:
            r_elem.find("a:t", nsmap).text = line
        else:
            r_elem = etree.SubElement(new_p, f'{{{nsmap["a"]}}}r')
            t_elem = etree.SubElement(r_elem, f'{{{nsmap["a"]}}}t')
            t_elem.text = line
        new_p.append(r_elem)
        txBody.append(new_p)


def update_table_cell(table, row, col, text):
    """Update text in a table cell."""
    cell = table.cell(row, col)
    if cell.text_frame.paragraphs[0].runs:
        cell.text_frame.paragraphs[0].runs[0].text = text
    else:
        cell.text = text


def main():
    prs = Presentation("Gaia_Presentation.pptx")

    # ================================================================
    # Slide 11 (index 10): Pkg 1 + Pkg 2
    # ================================================================
    slide11 = prs.slides[10]
    code_block = slide11.shapes[3]
    clear_and_set_text(code_block.text_frame, [
        "$ gaia init galileo_tied_balls",
        "",
        "# ═══ Pkg 1: aristotle_physics ═══",
        '$ gaia node add "越重的物体下落越快" \\',
        "    --prior 0.70 --type paper-extract   → 5001",
        '$ gaia node add "石头比树叶落得快" \\',
        "    --prior 0.90                        → 5002",
        "$ gaia edge add --tail 5001,5002 \\",
        "    --head 5003 --type abstraction",
        '  → 5003 "v ∝ W" prior=0.70',
        "$ gaia commit -m \"aristotle_physics\"",
        "$ gaia propagate",
        "  5003 (v ∝ W): belief = 0.70",
        "",
        "# ═══ Pkg 2: galileo1638_tied_balls ═══",
        '$ gaia node add "重球H绑轻球L"        → 5004',
        "$ gaia edge add --tail 5003,5004 --head 5005 \\",
        '    --type deduction  → "推导A: HL更慢"',
        "$ gaia edge add --tail 5003,5004 --head 5006 \\",
        '    --type deduction  → "推导B: HL更快"',
        "$ gaia edge add --tail 5005,5006 \\",
        "    --type contradiction",
        "  # 同一物体不可能既更快又更慢!",
        "$ gaia commit && gaia propagate",
        "  5003 (v∝W): 0.70→0.35 ↓  矛盾回传",
        "  5008 (v∝W错误): belief=0.82 ↑",
    ])

    # ================================================================
    # Slide 12 (index 11): Pkg 3 + Pkg 4 + Pkg 5
    # ================================================================
    slide12 = prs.slides[11]
    code_block = slide12.shapes[3]
    clear_and_set_text(code_block.text_frame, [
        "# ═══ Pkg 3: medium_density ═══",
        '$ gaia node add "水中落速差异更大" --prior 0.90',
        '$ gaia node add "介质密度↓ 差异↓" --prior 0.85',
        "$ gaia edge add --tail 5009,5010 --head 5011 \\",
        "    --type deduction",
        '  → 5011 "★ 结论①: 空气阻力才是混淆因素"',
        "$ gaia commit && gaia propagate",
        "  5003: 0.35→0.28 ↓  替代解释削弱旧定律",
        "",
        "# ═══ Pkg 4: vacuum_prediction ═══",
        "$ gaia edge add --tail 5008,5011 --head 5012",
        '  → 5012 "★ 结论②: 真空中所有物体等速下落"',
        '$ gaia node add "斜面实验: 不同球同时到底" \\',
        "    --prior 0.90",
        "$ gaia commit && gaia propagate",
        "  5012 (真空等速): 0.00→0.78 ↑",
        "",
        "# ═══ Pkg 5: newton_principia ═══",
        '$ gaia node add "F=ma" "F=mg" --prior 0.95',
        "$ gaia edge add --tail 5015,5016 --head 5017",
        '  → 5017 "a=g 质量约掉了!"',
        "$ gaia edge add --tail 5003,5017 \\",
        "    --type contradiction  # 第二条矛盾线!",
        "$ gaia commit && gaia propagate",
        "  5003→0.12↓↓  5017→0.93↑  5012→0.87↑",
    ])

    # Update diagram labels on slide 12
    for shape in slide12.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text
            if "0.35 → 0.08" in text:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.text = run.text.replace("0.35 → 0.08", "0.35 → 0.12")
            if "0.85 → 0.88" in text:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.text = run.text.replace("0.85 → 0.88", "0.78 → 0.87")
            if "从 0.35 降至 0.08" in text:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.text = run.text.replace("从 0.35 降至 0.08", "从 0.35 降至 0.12")

    # ================================================================
    # Slide 13 (index 12): Pkg 6 + Summary
    # ================================================================
    slide13 = prs.slides[12]
    code_block = slide13.shapes[3]
    clear_and_set_text(code_block.text_frame, [
        "# ═══ Pkg 6: apollo15_feather_drop ═══",
        '$ gaia node add "月球: ≈3e-12 atm" \\',
        "    --prior 0.99                      → 5018",
        '$ gaia node add "锤子=羽毛 同时落地" \\',
        "    --prior 0.99  # 质量比 44:1!      → 5019",
        "$ gaia edge add --tail 5018,5019 \\",
        "    --head 5020 --type deduction",
        "",
        "$ gaia commit -m \"apollo15\"",
        "$ gaia propagate",
        "  5003 (v∝W): 0.12→0.05 ↓  几乎归零",
        "  5012 (真空等速): 0.87→0.95 ↑  三线汇聚",
        "  5017 (a=g): 0.93→0.96 ↑  理论+实验",
        "  5020 (月球): belief=0.98 ↑  决定性",
        "",
        "$ gaia test     # 验证 beliefs.yaml",
        "$ gaia publish  # → 远程 registry",
    ])

    # Update table on slide 13
    table_shape = slide13.shapes[4]
    if table_shape.has_table:
        table = table_shape.table
        update_table_cell(table, 3, 1, "0.28 ↓")
        update_table_cell(table, 3, 2, "0.78")
        update_table_cell(table, 4, 1, "0.12 ↓")

    # ================================================================
    # Also update Einstein slides (15-17) to use CLI commands
    # ================================================================

    # Slide 15 (index 14): Einstein Pkg 1-3
    slide15 = prs.slides[14]
    code_block = slide15.shapes[3]
    clear_and_set_text(code_block.text_frame, [
        "$ gaia init einstein_elevator",
        "",
        "# ═══ Pkg 1: prior_knowledge ═══",
        '$ gaia node add "F=mᵢa" --prior 0.95        → 6001',
        '$ gaia node add "F=GMm/r²" --prior 0.95     → 6002',
        '$ gaia node add "光的微粒说" --prior 0.50     → 6003',
        '$ gaia node add "Maxwell EM" --prior 0.90    → 6005',
        '$ gaia node add "Eötvös mᵢ=mᵍ" --prior 0.95 → 6006',
        "$ gaia edge add --tail 6001,6002,6003 \\",
        '    --head 6004 --type deduction  P=0.60',
        '  → 6004 "Soldner: 0.87″" prior=0.60',
        "",
        "# ═══ Pkg 2: equivalence_principle ═══",
        '$ gaia node add "电梯思想实验" --prior 1.0  → 6007',
        "$ gaia edge add --tail 6006,6007 --head 6008",
        '  → 6008 "等效原理" belief≈0.85',
        "$ gaia edge add --tail 6008,6005 --head 6009",
        '  → 6009 "光在引力场弯曲"',
        "",
        "# ═══ Pkg 3: 1911_light_deflection ═══",
        "$ gaia edge add --tail 6009 --head 6010",
        '  → 6010 "Einstein: 0.87″" belief≈0.80',
        "$ gaia commit && gaia propagate",
        "# ⚠ 6004 = 6010 = 0.87″  此时无矛盾!",
    ])

    # Slide 16 (index 15): Einstein Pkg 4
    slide16 = prs.slides[15]
    code_block = slide16.shapes[3]
    clear_and_set_text(code_block.text_frame, [
        "# ═══ Pkg 4: general_relativity ═══",
        '$ gaia node add "GR: Gμν=8πGTμν" \\',
        "    --prior 0.85                      → 6012",
        "$ gaia edge add --tail 6012,6008 --head 6013",
        '  → 6013 "光沿零测地线 (时间+空间弯曲)"',
        "$ gaia edge add --tail 6013 --head 6014",
        '  → 6014 "GR: 1.75″ = 2×0.87″"',
        '$ gaia node add "水星进动 43″/世纪" \\',
        "    --prior 0.90                      → 6015",
        "$ gaia edge add --tail 6012 --head 6015",
        "",
        "# 关键: GR vs Newton 定量矛盾!",
        "$ gaia edge add --tail 6014,6004 \\",
        "    --type contradiction",
        "  # 1.75″ vs 0.87″ — 精确 2× 差异",
        "",
        "$ gaia commit -m \"general_relativity\"",
        "$ gaia propagate",
        "  6014 (GR 1.75″): belief=0.85",
        "  6004 (Newton 0.87″): 0.60→0.40 ↓",
    ])

    # Slide 17 (index 16): Einstein Pkg 5
    slide17 = prs.slides[16]
    code_block = slide17.shapes[3]
    clear_and_set_text(code_block.text_frame, [
        "# ═══ Pkg 5: eddington1919_eclipse ═══",
        '$ gaia node add "日食远征(Sobral+Príncipe)" \\',
        "    --prior 0.95                      → 6017",
        '$ gaia node add "观测 1.61±0.30″ / 1.98±0.16″" \\',
        "    --prior 0.90                      → 6018",
        "$ gaia edge add --tail 6017,6018 --head 6019",
        '  → 6019 "观测支持 GR, 排除 Newton"',
        "",
        "# 观测 vs 牛顿: 实验排除",
        "$ gaia edge add --tail 6019,6004 \\",
        "    --type contradiction  # ~1.7″ vs 0.87″",
        "",
        "$ gaia commit -m \"eddington1919\"",
        "$ gaia publish  # → 远程 registry",
        "$ gaia review <commit_id>",
        "$ gaia merge <commit_id>",
        "  6004: 0.40→0.10 ↓↓  6014: →0.95 ↑↑",
        "  6002 (Newton引力): 0.95→0.80 (降级)",
    ])

    # ================================================================
    prs.save("Gaia_Presentation.pptx")
    print("Updated all slides with gaia CLI commands")


if __name__ == "__main__":
    main()
