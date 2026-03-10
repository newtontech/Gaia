#!/usr/bin/env python3
"""Generate a ~25-slide PPTX presentation about the Gaia project.

Usage:
    python scripts/generate_presentation.py

Output:
    Gaia_Presentation.pptx in the project root directory.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Cm, Inches, Pt

# ── Colors ──────────────────────────────────────────────────────────────
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1A, 0x1A, 0x1A)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MID_GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xE0, 0xE0, 0xE0)
ACCENT_BLUE = RGBColor(0x1B, 0x5E, 0x9F)
ACCENT_LIGHT = RGBColor(0x42, 0xA5, 0xF5)
RED = RGBColor(0xD3, 0x2F, 0x2F)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
ORANGE = RGBColor(0xEF, 0x6C, 0x00)
TABLE_HEADER_BG = RGBColor(0x1B, 0x5E, 0x9F)
TABLE_ALT_BG = RGBColor(0xF5, 0xF5, 0xF5)

# Code box colors
CODE_BG = RGBColor(0x1E, 0x1E, 0x1E)
CODE_BORDER = RGBColor(0x3E, 0x3E, 0x3E)
CODE_COMMENT = RGBColor(0x6A, 0x9F, 0x55)
CODE_DEFAULT = RGBColor(0xD4, 0xD4, 0xD4)
CODE_KEYWORD = RGBColor(0x56, 0x9C, 0xD6)
CODE_WARN = RGBColor(0xFF, 0xD7, 0x00)
CODE_RED = RGBColor(0xF4, 0x47, 0x47)

# ── Fonts ───────────────────────────────────────────────────────────────
FONT_TITLE = "Helvetica Neue"
FONT_BODY = "Helvetica Neue"
FONT_MONO = "Menlo"

# ── Slide dimensions (16:9) ────────────────────────────────────────────
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def set_slide_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(
    slide,
    left,
    top,
    width,
    height,
    text,
    font_size=18,
    color=DARK_GRAY,
    bold=False,
    font_name=FONT_BODY,
    alignment=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    tf.vertical_anchor = anchor
    return txBox


def add_bullet_slide_text(
    tf, items, font_size=20, color=DARK_GRAY, bullet_color=ACCENT_BLUE, spacing=Pt(8)
):
    """Add bulleted items to an existing text frame."""
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = FONT_BODY
        p.level = 0
        p.space_after = spacing
        # Bullet character
        pPr = p._pPr
        if pPr is None:
            from pptx.oxml.ns import qn

            pPr = p._p.get_or_add_pPr()
        from pptx.oxml.ns import qn

        buNone = pPr.findall(qn("a:buNone"))
        for bn in buNone:
            pPr.remove(bn)
        # Add bullet
        from lxml import etree

        buChar = etree.SubElement(pPr, qn("a:buChar"))
        buChar.set("char", "●")
        buClr = etree.SubElement(pPr, qn("a:buClr"))
        srgbClr = etree.SubElement(buClr, qn("a:srgbClr"))
        srgbClr.set("val", str(bullet_color))


def add_section_header(slide, section_tag, slide_title):
    """Standard slide header with section indicator and title."""
    # Section tag
    add_text_box(
        slide,
        Cm(2),
        Cm(0.5),
        Cm(10),
        Cm(1),
        section_tag,
        font_size=12,
        color=ACCENT_BLUE,
        bold=True,
    )
    # Title
    add_text_box(
        slide, Cm(2), Cm(1.3), Cm(29), Cm(2), slide_title, font_size=32, color=BLACK, bold=True
    )
    # Divider line
    line = slide.shapes.add_connector(1, Cm(2), Cm(3.5), Cm(31), Cm(3.5))
    line.line.color.rgb = LIGHT_GRAY
    line.line.width = Pt(1)


def make_table(slide, left, top, width, rows_data, col_widths=None, font_size=14):
    """Create a styled table. rows_data[0] is the header row."""
    n_rows = len(rows_data)
    n_cols = len(rows_data[0])
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, Cm(1))
    table = table_shape.table

    # Set column widths
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    # Auto-size height
    Cm(0)
    for r_idx, row_data in enumerate(rows_data):
        row = table.rows[r_idx]
        row.height = Cm(1.2) if r_idx == 0 else Cm(1.0)

        for c_idx, cell_text in enumerate(row_data):
            cell = row.cells[c_idx]
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.text = str(cell_text)
            p.font.size = Pt(font_size)
            p.font.name = FONT_BODY

            if r_idx == 0:
                # Header
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_HEADER_BG
                p.font.color.rgb = WHITE
                p.font.bold = True
            else:
                if r_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = TABLE_ALT_BG
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = WHITE
                p.font.color.rgb = DARK_GRAY

            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p.alignment = PP_ALIGN.LEFT
            # Margins
            cell.margin_left = Cm(0.3)
            cell.margin_right = Cm(0.3)
            cell.margin_top = Cm(0.1)
            cell.margin_bottom = Cm(0.1)

    return table_shape


def draw_node(slide, cx, cy, w, h, text, color=ACCENT_BLUE, font_size=11, text_color=WHITE):
    """Draw a rounded rectangle node."""
    from pptx.enum.shapes import MSO_SHAPE

    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, cy, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.name = FONT_BODY
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    # Smaller margins
    tf.margin_left = Cm(0.15)
    tf.margin_right = Cm(0.15)
    tf.margin_top = Cm(0.05)
    tf.margin_bottom = Cm(0.05)
    return shape


def draw_arrow(slide, x1, y1, x2, y2, color=MID_GRAY, width=Pt(2)):
    """Draw a connector arrow."""
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)
    connector.line.color.rgb = color
    connector.line.width = width
    # Add arrowhead
    connector.line._ln.set("{http://schemas.openxmlformats.org/drawingml/2006/main}tailEnd", "")
    from lxml import etree
    from pptx.oxml.ns import qn

    tailEnd = etree.SubElement(connector.line._ln, qn("a:tailEnd"))
    tailEnd.set("type", "triangle")
    tailEnd.set("w", "med")
    tailEnd.set("len", "med")
    return connector


def draw_contradiction(slide, x1, y1, x2, y2):
    """Draw a red dashed contradiction line."""
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)
    connector.line.color.rgb = RED
    connector.line.width = Pt(2.5)
    connector.line.dash_style = 2  # dash
    return connector


def add_command_box(slide, left, top, width, height, lines):
    """Draw a dark code box with colored command lines.

    lines: list of (text, color) tuples.
    """
    from pptx.enum.shapes import MSO_SHAPE

    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = CODE_BG
    box.line.color.rgb = CODE_BORDER
    box.line.width = Pt(1)

    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Cm(0.4)
    tf.margin_right = Cm(0.3)
    tf.margin_top = Cm(0.3)
    tf.margin_bottom = Cm(0.3)

    for i, (text, color) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(11)
        p.font.name = FONT_BODY
        p.font.color.rgb = color
        p.space_after = Pt(1)

    return box


# ═══════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ═══════════════════════════════════════════════════════════════════════


def slide_01_cover(prs):
    """Slide 1: Title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, WHITE)

    # Accent bar at top
    from pptx.enum.shapes import MSO_SHAPE

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Cm(0.6))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_BLUE
    bar.line.fill.background()

    # Title
    add_text_box(
        slide,
        Cm(3),
        Cm(5),
        Cm(28),
        Cm(4),
        "Gaia",
        font_size=60,
        color=BLACK,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )

    add_text_box(
        slide,
        Cm(3),
        Cm(8.5),
        Cm(28),
        Cm(3),
        "面向 Agentic Science at Scale 的\n知识包管理系统与 Large Knowledge Model",
        font_size=26,
        color=MID_GRAY,
        alignment=PP_ALIGN.CENTER,
    )

    add_text_box(
        slide,
        Cm(3),
        Cm(13),
        Cm(28),
        Cm(2),
        "[演讲者]  ·  [日期]",
        font_size=18,
        color=MID_GRAY,
        alignment=PP_ALIGN.CENTER,
    )


def slide_02_agent_science(prs):
    """Slide 2: AI Agent as researcher."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_header(slide, "Part 1 · 为什么", "AI Agent 正在成为科研参与者")

    txBox = slide.shapes.add_textbox(Cm(2), Cm(4.5), Cm(18), Cm(12))
    tf = txBox.text_frame
    tf.word_wrap = True
    add_bullet_slide_text(
        tf,
        [
            "AlphaFold 2: 预测 2 亿蛋白质结构 → Nature 方法学年度突破",
            "GNoME: 发现 220 万新晶体结构，通过实验验证",
            "AlphaProof / AlphaGeometry: IMO 级别数学推理",
            "自动化实验室：Agent 设计-执行-分析闭环",
            "趋势：Agent 不再只是工具，而是 co-researcher",
        ],
        font_size=20,
    )

    # Right side callout
    from pptx.enum.shapes import MSO_SHAPE

    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(22), Cm(5), Cm(10), Cm(8))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xE3, 0xF2, 0xFD)
    box.line.color.rgb = ACCENT_BLUE
    box.line.width = Pt(1.5)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Cm(0.5)
    tf.margin_right = Cm(0.5)
    tf.margin_top = Cm(0.5)
    p = tf.paragraphs[0]
    p.text = "但两个瓶颈正在显现 ↓"
    p.font.size = Pt(22)
    p.font.color.rgb = ACCENT_BLUE
    p.font.bold = True
    p.font.name = FONT_BODY
    p.alignment = PP_ALIGN.CENTER


def slide_03_bottleneck1(prs):
    """Slide 3: Bottleneck 1 — No reliable knowledge foundation."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_header(slide, "Part 1 · 为什么", "瓶颈 1 — Agent 缺乏可靠的知识底座")

    txBox = slide.shapes.add_textbox(Cm(2), Cm(4.5), Cm(18), Cm(12))
    tf = txBox.text_frame
    tf.word_wrap = True
    add_bullet_slide_text(
        tf,
        [
            "论文间矛盾无法系统追踪",
            "撤稿 (retraction) 的连锁影响无法传播",
            "结论的可信度无法量化",
            "引用 ≠ 认同，但现有系统无法区分",
        ],
        font_size=22,
    )

    # Diagram: retracted paper impact
    draw_node(slide, Cm(22), Cm(5), Cm(4.5), Cm(1.6), "Paper A", ACCENT_BLUE, 12)
    draw_node(slide, Cm(22), Cm(8), Cm(4.5), Cm(1.6), "Paper B\n(Retracted!)", RED, 12)
    draw_node(slide, Cm(22), Cm(11), Cm(4.5), Cm(1.6), "Paper C", ACCENT_BLUE, 12)
    draw_node(slide, Cm(27.5), Cm(6.5), Cm(4.5), Cm(1.6), "Paper D", MID_GRAY, 12)
    draw_node(slide, Cm(27.5), Cm(9.5), Cm(4.5), Cm(1.6), "Paper E", MID_GRAY, 12)

    # Arrows
    draw_arrow(slide, Cm(24.25), Cm(6.6), Cm(24.25), Cm(8), MID_GRAY)
    draw_arrow(slide, Cm(24.25), Cm(9.6), Cm(24.25), Cm(11), MID_GRAY)
    draw_arrow(slide, Cm(26.5), Cm(8.5), Cm(27.5), Cm(7.8), MID_GRAY)
    draw_arrow(slide, Cm(26.5), Cm(9), Cm(27.5), Cm(9.8), MID_GRAY)

    add_text_box(
        slide,
        Cm(22),
        Cm(13.5),
        Cm(10),
        Cm(1.5),
        "↑ 撤稿影响无法自动传播到 D, E",
        font_size=13,
        color=RED,
        alignment=PP_ALIGN.CENTER,
    )


def slide_04_bottleneck2(prs):
    """Slide 4: Bottleneck 2 — Peer review can't scale."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_header(slide, "Part 1 · 为什么", "瓶颈 2 — 同行评审无法 Scale")

    txBox = slide.shapes.add_textbox(Cm(2), Cm(4.5), Cm(20), Cm(12))
    tf = txBox.text_frame
    tf.word_wrap = True
    add_bullet_slide_text(
        tf,
        [
            "Agent 产出速度 >> 人类审稿能力",
            "当前同行评审：周期数月、人力瓶颈",
            "需要：机器可操作的知识管理与质量保障机制",
            "不是替代人类审稿，而是让机器+人类协同审稿",
        ],
        font_size=22,
    )

    # Callout: introducing Gaia
    from pptx.enum.shapes import MSO_SHAPE

    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(22), Cm(7), Cm(10), Cm(5))
    box.fill.solid()
    box.fill.fore_color.rgb = ACCENT_BLUE
    box.line.fill.background()
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "我们的方案：\nGaia"
    p.font.size = Pt(28)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_BODY
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE


def slide_05_core_idea(prs):
    """Slide 5: Gaia core idea."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_header(slide, "Part 2 · Gaia 是什么", "核心思想 — 知识 = 概率推理图")

    txBox = slide.shapes.add_textbox(Cm(2), Cm(4.5), Cm(14), Cm(10))
    tf = txBox.text_frame
    tf.word_wrap = True
    add_bullet_slide_text(
        tf,
        [
            "每条知识 = 一个命题节点 (Node)，带可信度 (belief)",
            "推理关系 = 超边 (HyperEdge)，连接前提与结论",
            "新证据加入 → Belief Propagation 沿推理链自动传播",
            "矛盾 = 一等公民，用 contradiction edge 显式建模",
        ],
        font_size=20,
    )

    # Simple 3-node diagram
    draw_node(slide, Cm(19), Cm(5), Cm(5.5), Cm(2), "前提 A\nbelief=0.9", ACCENT_BLUE, 13)
    draw_node(slide, Cm(26), Cm(5), Cm(5.5), Cm(2), "前提 B\nbelief=0.8", ACCENT_BLUE, 13)
    draw_node(slide, Cm(22.5), Cm(11), Cm(5.5), Cm(2), "结论 C\nbelief=?", GREEN, 13)

    # Reasoning edge box
    from pptx.enum.shapes import MSO_SHAPE

    ebox = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(22), Cm(8.2), Cm(6.5), Cm(1.5))
    ebox.fill.solid()
    ebox.fill.fore_color.rgb = ORANGE
    ebox.line.fill.background()
    tf = ebox.text_frame
    p = tf.paragraphs[0]
    p.text = "HyperEdge P=0.85"
    p.font.size = Pt(12)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    draw_arrow(slide, Cm(21.75), Cm(7), Cm(23.5), Cm(8.2), MID_GRAY)
    draw_arrow(slide, Cm(28.75), Cm(7), Cm(27), Cm(8.2), MID_GRAY)
    draw_arrow(slide, Cm(25.25), Cm(9.7), Cm(25.25), Cm(11), MID_GRAY)


def slide_06_paradigm_shift(prs):
    """Slide 6: Proposition-level vs Entity-level."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_header(slide, "Part 2 · Gaia 是什么", "命题级 vs 实体级 — 范式转换")

    rows = [
        ["", "传统 Knowledge Graph", "Gaia"],
        ["基本单元", "Entity", "Proposition (命题)"],
        ["关系", "Entity → Relation → Entity", "Proposition → Reasoning → Proposition"],
        ["示例", '"Einstein bornIn Ulm"', '"Einstein was born in Ulm" belief=0.99'],
        ["真值", "存了就是真的 (binary)", "每条知识带可信度 (probabilistic)"],
        ["矛盾处理", "无法表示", "Contradiction Edge (一等公民)"],
        ["推理传播", "无", "Belief Propagation 自动传播"],
    ]

    make_table(
        slide, Cm(2), Cm(4.5), Cm(30), rows, col_widths=[Cm(4.5), Cm(12), Cm(13.5)], font_size=15
    )


def slide_07_contradiction(prs):
    """Slide 7: Contradiction as first-class citizen."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_header(slide, "Part 2 · Gaia 是什么", "Contradiction Edge — 矛盾是一等公民")

    txBox = slide.shapes.add_textbox(Cm(2), Cm(4.5), Cm(16), Cm(10))
    tf = txBox.text_frame
    tf.word_wrap = True
    add_bullet_slide_text(
        tf,
        [
            "在 Gaia 中，矛盾不是 bug，是 feature",
            "两个理论做出不同预测 → contradiction edge 连接",
            "Belief Propagation 根据证据自动分配可信度",
            "旧理论不删除，belief 下降 — 保留历史完整性",
        ],
        font_size=20,
    )

    # A ↔ B contradiction diagram
    draw_node(slide, Cm(21), Cm(5.5), Cm(5), Cm(2.5), '理论 A\n预测: 1.75"', ACCENT_BLUE, 14)
    draw_node(slide, Cm(28), Cm(5.5), Cm(5), Cm(2.5), '理论 B\n预测: 0.87"', ACCENT_BLUE, 14)

    # Contradiction line
    draw_contradiction(slide, Cm(26), Cm(6.75), Cm(28), Cm(6.75))

    add_text_box(
        slide,
        Cm(23.5),
        Cm(8.5),
        Cm(7),
        Cm(1.5),
        "← Contradiction Edge →",
        font_size=14,
        color=RED,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )

    # Result
    draw_node(slide, Cm(21), Cm(11), Cm(5), Cm(2), "belief ↑ 0.85", GREEN, 14)
    draw_node(slide, Cm(28), Cm(11), Cm(5), Cm(2), "belief ↓ 0.15", RED, 14, WHITE)

    add_text_box(
        slide,
        Cm(23.5),
        Cm(10),
        Cm(7),
        Cm(1),
        "观测支持 A 后：",
        font_size=14,
        color=MID_GRAY,
        alignment=PP_ALIGN.CENTER,
    )


def slide_08_bp(prs):
    """Slide 8: Factor Graph + BP."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_header(slide, "Part 2 · Gaia 是什么", "Factor Graph + Belief Propagation")

    # Formula
    add_text_box(
        slide,
        Cm(2),
        Cm(4.5),
        Cm(30),
        Cm(2),
        "belief(结论) ∝ ∏ belief(前提ᵢ) × P(推理)",
        font_size=28,
        color=BLACK,
        bold=True,
        font_name=FONT_MONO,
        alignment=PP_ALIGN.CENTER,
    )

    # Intuition
    txBox = slide.shapes.add_textbox(Cm(2), Cm(7), Cm(15), Cm(8))
    tf = txBox.text_frame
    tf.word_wrap = True
    add_bullet_slide_text(
        tf,
        [
            "消息传递：前提可信度 × 推理可靠度 → 结论可信度",
            "多条证据汇聚 → 结论更可靠 (corroboration)",
            "矛盾证据 → 双方 belief 下降 (competition)",
            "Loopy BP：大规模图上近似推断",
        ],
        font_size=20,
    )

    # Message passing diagram
    draw_node(slide, Cm(20), Cm(7.5), Cm(4), Cm(1.6), "前提₁ 0.9", ACCENT_BLUE, 12)
    draw_node(slide, Cm(20), Cm(10), Cm(4), Cm(1.6), "前提₂ 0.8", ACCENT_BLUE, 12)

    from pptx.enum.shapes import MSO_SHAPE

    f = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(25.5), Cm(8.5), Cm(2.5), Cm(2))
    f.fill.solid()
    f.fill.fore_color.rgb = ORANGE
    f.line.fill.background()
    tf = f.text_frame
    p = tf.paragraphs[0]
    p.text = "f\nP=0.9"
    p.font.size = Pt(12)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    draw_node(slide, Cm(29.5), Cm(8.5), Cm(3.5), Cm(2), "结论\n0.65", GREEN, 13)

    draw_arrow(slide, Cm(24), Cm(8.3), Cm(25.5), Cm(9), ACCENT_BLUE)
    draw_arrow(slide, Cm(24), Cm(10.8), Cm(25.5), Cm(10), ACCENT_BLUE)
    draw_arrow(slide, Cm(28), Cm(9.5), Cm(29.5), Cm(9.5), ORANGE, Pt(2.5))

    add_text_box(
        slide,
        Cm(24.5),
        Cm(12),
        Cm(4),
        Cm(1),
        "msg →",
        font_size=11,
        color=ORANGE,
        alignment=PP_ALIGN.CENTER,
    )


def slide_09_git_workflow(prs):
    """Slide 9: Git-like commit workflow."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_header(slide, "Part 2 · Gaia 是什么", "Git-like 提交工作流")

    txBox = slide.shapes.add_textbox(Cm(2), Cm(4.5), Cm(15), Cm(10))
    tf = txBox.text_frame
    tf.word_wrap = True
    add_bullet_slide_text(
        tf,
        [
            "Submit: 提交新知识（节点 + 超边）",
            "Review: 自动/人工质量检查",
            "Merge: 写入存储 + 触发 BP 传播",
            "每一步都有质量门控",
        ],
        font_size=20,
    )

    # Pipeline diagram
    from pptx.enum.shapes import MSO_SHAPE

    steps = [("Submit", ACCENT_BLUE), ("Review", ORANGE), ("Merge", GREEN)]
    x_start = Cm(19)
    for i, (label, color) in enumerate(steps):
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x_start + Cm(i * 5), Cm(6), Cm(4), Cm(2.5)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(20)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        if i < 2:
            draw_arrow(
                slide,
                x_start + Cm(i * 5 + 4),
                Cm(7.25),
                x_start + Cm((i + 1) * 5),
                Cm(7.25),
                MID_GRAY,
                Pt(3),
            )

    # Sub-labels
    sub_labels = ["AddNode\nAddEdge", "LLM 审查\n一致性检查", "Triple-Write\nBP 传播"]
    for i, label in enumerate(sub_labels):
        add_text_box(
            slide,
            x_start + Cm(i * 5),
            Cm(9),
            Cm(4),
            Cm(2),
            label,
            font_size=12,
            color=MID_GRAY,
            alignment=PP_ALIGN.CENTER,
        )


def slide_10_galileo_overview(prs):
    """Slide 10: Galileo overview."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_header(slide, "Part 3 · 伽利略思想实验", "伽利略绑球实验 — 概述")

    txBox = slide.shapes.add_textbox(Cm(2), Cm(4.5), Cm(20), Cm(10))
    tf = txBox.text_frame
    tf.word_wrap = True
    add_bullet_slide_text(
        tf,
        [
            "1638 年《Discorsi》: 一个思想实验推翻 2000 年学说",
            "6 个 knowledge package，跨越 2300 年",
            "演示 Gaia 的核心机制：contradiction edge 驱动 belief 修正",
            "全程只添加 node 和 edge，从不修改或删除历史",
        ],
        font_size=22,
    )

    # Timeline
    from pptx.enum.shapes import MSO_SHAPE

    line = slide.shapes.add_connector(1, Cm(2), Cm(14.5), Cm(31), Cm(14.5))
    line.line.color.rgb = ACCENT_BLUE
    line.line.width = Pt(2)

    pkgs = [
        ("Pkg 1", "350 BCE\n亚里士多德"),
        ("Pkg 2", "1638\n绑球悖论"),
        ("Pkg 3-4", "1638\n密度+真空"),
        ("Pkg 5", "1687\n牛顿"),
        ("Pkg 6", "1971\nApollo 15"),
    ]
    for i, (tag, label) in enumerate(pkgs):
        x = Cm(2.5 + i * 6)
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Cm(0.7), Cm(14), Cm(0.9), Cm(0.9))
        dot.fill.solid()
        dot.fill.fore_color.rgb = ACCENT_BLUE
        dot.line.fill.background()
        add_text_box(
            slide,
            x - Cm(0.5),
            Cm(15.5),
            Cm(4),
            Cm(2.5),
            f"{tag}\n{label}",
            font_size=10,
            color=MID_GRAY,
            alignment=PP_ALIGN.CENTER,
        )

    # Right: belief trajectory preview
    add_text_box(
        slide,
        Cm(22),
        Cm(5),
        Cm(10),
        Cm(1.5),
        "v ∝ W belief 轨迹",
        font_size=16,
        color=ACCENT_BLUE,
        bold=True,
    )
    trajectory = [
        ("Pkg 1  先验", "0.70", MID_GRAY),
        ("Pkg 2  绑球矛盾", "0.35 ↓", ORANGE),
        ("Pkg 3-4 密度+真空", "0.15 ↓", ORANGE),
        ("Pkg 5  Newton a=g", "0.08 ↓", RED),
        ("Pkg 6  Apollo 15", "0.05 ↓", RED),
    ]
    for i, (stage, belief, color) in enumerate(trajectory):
        y = Cm(6.5 + i * 1.5)
        add_text_box(slide, Cm(22), y, Cm(6), Cm(1.2), stage, font_size=13, color=DARK_GRAY)
        add_text_box(slide, Cm(29), y, Cm(3), Cm(1.2), belief, font_size=14, color=color, bold=True)


def slide_11_galileo_pkg12(prs):
    """Slide 11: Galileo Pkg 1+2 — Aristotle + Tied Balls contradiction."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_header(slide, "Part 3 · 伽利略思想实验", "Pkg 1-2: 亚里士多德先验 + 绑球悖论")

    # ── Left: command box ──
    add_command_box(
        slide,
        Cm(1.5),
        Cm(4.2),
        Cm(13.5),
        Cm(14),
        [
            ("# Pkg 1: aristotle_physics", CODE_COMMENT),
            ("POST /commits {", CODE_KEYWORD),
            (' message: "aristotle_physics",', CODE_DEFAULT),
            (" operations: [", CODE_DEFAULT),
            ('  {op:"add_edge", type:"abstraction",', CODE_DEFAULT),
            ('   tail: ["自然运动学说","石头快于叶子"],', CODE_DEFAULT),
            ('   head: ["v ∝ W 定律"]}', CODE_DEFAULT),
            (" ]}", CODE_DEFAULT),
            ("", CODE_DEFAULT),
            ("# Pkg 2: galileo1638_tied_balls", CODE_COMMENT),
            ("POST /commits {", CODE_KEYWORD),
            (' message: "galileo1638_tied_balls",', CODE_DEFAULT),
            (" operations: [", CODE_DEFAULT),
            ('  {op:"add_edge", type:"deduction",', CODE_DEFAULT),
            ('   tail: [#5003"v∝W", "绑球设定 H+L"],', CODE_DEFAULT),
            ('   head: ["推导A: L拖慢H → 更慢"]},', CODE_DEFAULT),
            ('  {op:"add_edge", type:"deduction",', CODE_DEFAULT),
            ("   tail: [#5003, #5004],", CODE_DEFAULT),
            ('   head: ["推导B: 组合更重 → 更快"]},', CODE_DEFAULT),
            ('  {op:"add_edge", type:"contradiction",', CODE_RED),
            ('   tail: [#5005"推导A", #5006"推导B"],', CODE_RED),
            ("   head: []}", CODE_RED),
            (" ]}", CODE_DEFAULT),
            ("# → BP: 5003 belief 0.70 → 0.35 ↓", CODE_WARN),
        ],
    )

    # ── Right: graph ──
    # Row 1: premises
    draw_node(slide, Cm(17), Cm(4.3), Cm(4.5), Cm(1.6), "5001\n自然运动", MID_GRAY, 10)
    draw_node(slide, Cm(22.5), Cm(4.3), Cm(4.5), Cm(1.6), "5002\n石头>叶子", MID_GRAY, 10)

    # Row 2: Aristotle's law + setup
    draw_node(slide, Cm(16), Cm(7.2), Cm(5.5), Cm(2), "5003 v ∝ W\n0.70 → 0.35 ↓", RED, 12, WHITE)
    draw_node(slide, Cm(23), Cm(7.2), Cm(4.5), Cm(2), "5004\n绑球设定", MID_GRAY, 11)

    # Arrows: 5001,5002 → 5003
    draw_arrow(slide, Cm(19.25), Cm(5.9), Cm(18.75), Cm(7.2), MID_GRAY)
    draw_arrow(slide, Cm(24.75), Cm(5.9), Cm(19.5), Cm(7.2), MID_GRAY)

    # Row 3: deductions A and B
    draw_node(slide, Cm(16), Cm(11), Cm(5), Cm(2), "5005\n推导A: 更慢", ACCENT_BLUE, 11)
    draw_node(slide, Cm(23), Cm(11), Cm(5), Cm(2), "5006\n推导B: 更快", ACCENT_BLUE, 11)

    # Arrows: 5003+5004 → 5005, 5003+5004 → 5006
    draw_arrow(slide, Cm(18.75), Cm(9.2), Cm(18.5), Cm(11), MID_GRAY)
    draw_arrow(slide, Cm(19.5), Cm(9.2), Cm(25.5), Cm(11), MID_GRAY)
    draw_arrow(slide, Cm(25.25), Cm(9.2), Cm(25.5), Cm(11), MID_GRAY)

    # Contradiction edge between 5005 and 5006
    draw_contradiction(slide, Cm(21), Cm(12), Cm(23), Cm(12))
    add_text_box(
        slide,
        Cm(20.5),
        Cm(13.2),
        Cm(3.5),
        Cm(1),
        "CONTRADICTION",
        font_size=10,
        color=RED,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )

    # Row 4: result
    draw_node(slide, Cm(18.5), Cm(15), Cm(8), Cm(2), "5008: v∝W 自相矛盾  belief=0.85", GREEN, 12)
    draw_arrow(slide, Cm(21), Cm(13.5), Cm(22.5), Cm(15), MID_GRAY)

    # Legend
    add_text_box(
        slide,
        Cm(28),
        Cm(7),
        Cm(5),
        Cm(5),
        "图例\n● Node\n— 推理边\n- - 矛盾边",
        font_size=10,
        color=MID_GRAY,
    )


def slide_12_galileo_pkg345(prs):
    """Slide 12: Galileo Pkg 3+4+5 — evidence accumulation."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_header(slide, "Part 3 · 伽利略思想实验", "Pkg 3-5: 密度观察 + 真空预测 + 牛顿推导")

    # ── Left: command box ──
    add_command_box(
        slide,
        Cm(1.5),
        Cm(4.2),
        Cm(13.5),
        Cm(14),
        [
            ("# Pkg 3: galileo1638_medium_density", CODE_COMMENT),
            ("POST /commits {", CODE_KEYWORD),
            (' operations: [{op:"add_edge",', CODE_DEFAULT),
            ('  type:"deduction",', CODE_DEFAULT),
            ('  tail:["密度递减观察","速差随密度↓"],', CODE_DEFAULT),
            ('  head:["空气阻力是混淆变量"]}]}', CODE_DEFAULT),
            ("", CODE_DEFAULT),
            ("# Pkg 4: galileo1638_vacuum_prediction", CODE_COMMENT),
            ("POST /commits {", CODE_KEYWORD),
            (' operations: [{op:"add_edge",', CODE_DEFAULT),
            ('  type:"deduction",', CODE_DEFAULT),
            ('  tail:[#5008"v∝W矛盾",#5011"空气阻力"],', CODE_DEFAULT),
            ('  head:["真空中所有物体等速下落"]}]}', CODE_DEFAULT),
            ("", CODE_DEFAULT),
            ("# Pkg 5: newton1687_principia", CODE_COMMENT),
            ("POST /commits {", CODE_KEYWORD),
            (" operations: [", CODE_DEFAULT),
            ('  {op:"add_edge", type:"deduction",', CODE_DEFAULT),
            ('   tail:["F=ma","F=mg"],', CODE_DEFAULT),
            ('   head:["∴ a=g 加速度与质量无关"]},', CODE_DEFAULT),
            ('  {op:"add_edge", type:"contradiction",', CODE_RED),
            ('   tail:[#5003"v∝W", #5017"a=g"],', CODE_RED),
            ("   head:[]}]}", CODE_RED),
            ("# → BP: 5003→0.08↓  5017→0.93↑", CODE_WARN),
        ],
    )

    # ── Right: graph — show key nodes and connections ──
    # Top: 5003 (Aristotle, dropping)
    draw_node(slide, Cm(17), Cm(4.5), Cm(5), Cm(1.8), "5003 v∝W\n0.35 → 0.08 ↓↓", RED, 11, WHITE)

    # Middle-left: 5011 (air resistance)
    draw_node(slide, Cm(16), Cm(8), Cm(5), Cm(1.8), "5011 空气阻力\nbelief=0.80", ACCENT_BLUE, 11)

    # Middle-right: 5017 (a=g)
    draw_node(slide, Cm(24), Cm(8), Cm(5), Cm(1.8), "5017 a=g\nbelief=0.93", GREEN, 11)

    # Contradiction between 5003 and 5017
    draw_contradiction(slide, Cm(22), Cm(5.4), Cm(25.5), Cm(8))
    add_text_box(
        slide, Cm(24), Cm(6), Cm(5), Cm(1), "CONTRADICTION②", font_size=10, color=RED, bold=True
    )

    # Arrow from 5015,5016 to 5017
    draw_node(slide, Cm(24), Cm(4.5), Cm(2), Cm(1.2), "F=ma", MID_GRAY, 9)
    draw_node(slide, Cm(27), Cm(4.5), Cm(2), Cm(1.2), "F=mg", MID_GRAY, 9)
    draw_arrow(slide, Cm(25), Cm(5.7), Cm(26), Cm(8), MID_GRAY)
    draw_arrow(slide, Cm(28), Cm(5.7), Cm(27), Cm(8), MID_GRAY)

    # Bottom: 5012 (vacuum prediction, rising)
    draw_node(slide, Cm(19), Cm(12), Cm(6.5), Cm(2), "5012 真空等速预测\n0.85 → 0.88 ↑", GREEN, 12)

    # Arrows: 5008+5011 → 5012
    draw_arrow(slide, Cm(18.5), Cm(9.8), Cm(21), Cm(12), MID_GRAY)
    # Arrow: 5017 → 5012 (independent confirmation)
    draw_arrow(slide, Cm(26.5), Cm(9.8), Cm(24), Cm(12), ACCENT_BLUE)

    # Annotation
    add_text_box(
        slide,
        Cm(26),
        Cm(12),
        Cm(7),
        Cm(2),
        "Newton 从 F=ma, F=mg\n独立推导出 a=g\n→ 第二条矛盾线",
        font_size=12,
        color=MID_GRAY,
    )

    # Bottom belief summary
    from pptx.enum.shapes import MSO_SHAPE

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(16), Cm(15.5), Cm(17), Cm(2.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
    bar.line.fill.background()
    add_text_box(
        slide,
        Cm(16.5),
        Cm(15.7),
        Cm(16),
        Cm(2),
        "两条独立矛盾线 + 空气阻力替代解释 → v∝W 从 0.35 降至 0.08",
        font_size=13,
        color=DARK_GRAY,
        bold=True,
    )


def slide_13_galileo_pkg6(prs):
    """Slide 13: Galileo Pkg 6 — Apollo + final beliefs."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_header(
        slide, "Part 3 · 伽利略思想实验", "Pkg 6: Apollo 15 — 决定性实验 + Belief 全景"
    )

    # ── Left: command box ──
    add_command_box(
        slide,
        Cm(1.5),
        Cm(4.2),
        Cm(13.5),
        Cm(5.5),
        [
            ("# Pkg 6: apollo15_1971_feather_drop", CODE_COMMENT),
            ("POST /commits {", CODE_KEYWORD),
            (" operations: [", CODE_DEFAULT),
            ('  {op:"add_edge", type:"deduction",', CODE_DEFAULT),
            ('   tail:["月球真空环境","锤=羽同时落"],', CODE_DEFAULT),
            ('   head:["等速落体实验确认"]},', CODE_DEFAULT),
            ('  {op:"add_edge", type:"deduction",', CODE_DEFAULT),
            ("   tail:[#5020], head:[#5012]}]}", CODE_DEFAULT),
        ],
    )

    # ── Left: belief evolution table ──
    rows = [
        ["阶段", "v∝W (旧)", "a=g (新)", "机制"],
        ["Pkg 1 先验", "0.70", "—", "日常观察"],
        ["Pkg 2 绑球", "0.35 ↓", "—", "Contradiction ①"],
        ["Pkg 3-4 密度+真空", "0.15 ↓", "0.85", "替代解释+实验"],
        ["Pkg 5 Newton", "0.08 ↓", "0.93 ↑", "理论 + Contradiction ②"],
        ["Pkg 6 Apollo", "0.05 ↓", "0.98 ↑", "月球真空实验"],
    ]
    make_table(
        slide,
        Cm(1.5),
        Cm(10.5),
        Cm(13.5),
        rows,
        col_widths=[Cm(4), Cm(2.5), Cm(2.5), Cm(4.5)],
        font_size=12,
    )

    # ── Right: convergence diagram ──
    add_text_box(
        slide,
        Cm(17),
        Cm(4.2),
        Cm(15),
        Cm(1.2),
        "三条独立证据线汇聚:",
        font_size=16,
        color=BLACK,
        bold=True,
    )

    from pptx.enum.shapes import MSO_SHAPE

    lines_data = [
        ("① 逻辑矛盾", "绑球悖论\nContradiction Edge", RED),
        ("② 理论推导", "Newton: F=ma, F=mg\n→ a=g + Contradiction", ACCENT_BLUE),
        ("③ 直接实验", "Apollo 15\n锤子 = 羽毛", GREEN),
    ]

    for i, (tag, desc, color) in enumerate(lines_data):
        y = Cm(5.8 + i * 3)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(17), y, Cm(3.5), Cm(2))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = tag
        p.font.size = Pt(13)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        add_text_box(
            slide,
            Cm(21),
            y,
            Cm(6),
            Cm(2),
            desc,
            font_size=11,
            color=DARK_GRAY,
            anchor=MSO_ANCHOR.MIDDLE,
        )

        draw_arrow(slide, Cm(27), y + Cm(1), Cm(28), y + Cm(1), color, Pt(2))

    # Convergence target
    draw_node(slide, Cm(28.5), Cm(7.5), Cm(4.5), Cm(4), "5012\n真空等速\nbelief\n0.98", GREEN, 13)

    # Bottom: key takeaway
    add_text_box(
        slide,
        Cm(17),
        Cm(15.5),
        Cm(16),
        Cm(2),
        "旧理论不删除，belief 下降；全程只添加 node+edge，从不修改历史",
        font_size=14,
        color=ACCENT_BLUE,
        bold=True,
    )


def slide_14_einstein_overview(prs):
    """Slide 14: Einstein elevator overview."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_header(slide, "Part 4 · 爱因斯坦思想实验", "爱因斯坦电梯 — 概述")

    txBox = slide.shapes.add_textbox(Cm(2), Cm(4.5), Cm(18), Cm(8))
    tf = txBox.text_frame
    tf.word_wrap = True
    add_bullet_slide_text(
        tf,
        [
            "5 个 package (1801–1919)：从牛顿到广义相对论",
            '核心叙事：0.87" → 0.87" → 1.75" → 观测 ~1.7"',
            "牛顿和爱因斯坦最初预测一致 → 无矛盾",
            "GR 预测翻倍 → 矛盾出现 → Eddington 判决",
        ],
        font_size=22,
    )

    # Visual: prediction journey
    from pptx.enum.shapes import MSO_SHAPE

    entries = [
        ("Pkg 1\n1801 Soldner", '0.87"', ACCENT_BLUE, "Newton 粒子引力"),
        ("Pkg 2-3\n1907-1911", '0.87"', ACCENT_BLUE, "等效原理\n(和 Newton 一样!)"),
        ("Pkg 4\n1915 GR", '1.75"', ORANGE, "时空弯曲\n(精确 2×!)"),
        ("Pkg 5\n1919 观测", '~1.7"', GREEN, "Eddington\n日食观测"),
    ]
    for i, (label, val, color, note) in enumerate(entries):
        x = Cm(2 + i * 8)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Cm(12.5), Cm(6.5), Cm(2.2))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"{label}  →  {val}"
        p.font.size = Pt(13)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        add_text_box(
            slide,
            x,
            Cm(15),
            Cm(6.5),
            Cm(2),
            note,
            font_size=11,
            color=MID_GRAY,
            alignment=PP_ALIGN.CENTER,
        )

        if i < 3:
            draw_arrow(slide, x + Cm(6.5), Cm(13.6), x + Cm(8), Cm(13.6), MID_GRAY, Pt(2))


def slide_15_einstein_pkg123(prs):
    """Slide 15: Einstein Pkg 1+2+3 — prior + equivalence + 1911."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_header(
        slide, "Part 4 · 爱因斯坦思想实验", "Pkg 1-3: 先验知识 → 等效原理 → 1911 预测"
    )

    # ── Left: command box ──
    add_command_box(
        slide,
        Cm(1.5),
        Cm(4.2),
        Cm(14),
        Cm(14),
        [
            ("# Pkg 1: prior_knowledge", CODE_COMMENT),
            ('AddNode 6001 "F = m_i·a"     prior=0.95', CODE_DEFAULT),
            ('AddNode 6002 "F = GMm/r²"    prior=0.95', CODE_DEFAULT),
            ('AddNode 6003 "光的微粒说"      prior=0.5', CODE_DEFAULT),
            ('AddNode 6004 "Soldner: 0.87″" prior=0.6', CODE_KEYWORD),
            ('AddNode 6005 "Maxwell EM"     prior=0.9', CODE_DEFAULT),
            ('AddNode 6006 "Eötvös m_i=m_g" prior=0.95', CODE_DEFAULT),
            ("AddEdge 6001 deduction", CODE_DEFAULT),
            ("  [6001,6002,6003] → [6004]  P=0.6", CODE_DEFAULT),
            ("", CODE_DEFAULT),
            ("# Pkg 2: einstein1907_equivalence", CODE_COMMENT),
            ('AddNode 6007 "电梯思想实验"    prior=1.0', CODE_DEFAULT),
            ('AddNode 6008 "等效原理"        prior=0.85', CODE_KEYWORD),
            ('AddNode 6009 "光在引力场弯曲"  prior=0.85', CODE_DEFAULT),
            ("AddEdge 6002 deduction", CODE_DEFAULT),
            ("  [6006,6007] → [6008]  P=0.85", CODE_DEFAULT),
            ("AddEdge 6003 deduction", CODE_DEFAULT),
            ("  [6008,6005] → [6009]  P=0.85", CODE_DEFAULT),
            ("", CODE_DEFAULT),
            ("# Pkg 3: einstein1911_light_deflection", CODE_COMMENT),
            ('AddNode 6010 "Einstein: 0.87″" prior=0.8', CODE_KEYWORD),
            ("AddEdge 6004 deduction", CODE_DEFAULT),
            ("  [6009] → [6010]  P=0.8", CODE_DEFAULT),
            ("", CODE_DEFAULT),
            ("# ⚠ 6004 = 6010 = 0.87″  此时无矛盾!", CODE_WARN),
        ],
    )

    # ── Right: graph ──
    # Newton path
    draw_node(slide, Cm(17), Cm(4.5), Cm(3.5), Cm(1.4), "6001\nF=m·a", MID_GRAY, 9)
    draw_node(slide, Cm(21), Cm(4.5), Cm(3.5), Cm(1.4), "6002\nF=GMm/r²", MID_GRAY, 9)
    draw_node(slide, Cm(25), Cm(4.5), Cm(3.5), Cm(1.4), "6003\n微粒说", MID_GRAY, 9)

    draw_node(
        slide, Cm(20.5), Cm(7.5), Cm(5), Cm(1.8), '6004 Soldner\n0.87" belief≈0.6', ACCENT_BLUE, 11
    )

    draw_arrow(slide, Cm(18.75), Cm(5.9), Cm(22), Cm(7.5), MID_GRAY)
    draw_arrow(slide, Cm(22.75), Cm(5.9), Cm(23), Cm(7.5), MID_GRAY)
    draw_arrow(slide, Cm(26.75), Cm(5.9), Cm(24), Cm(7.5), MID_GRAY)

    # Einstein path
    draw_node(slide, Cm(17), Cm(10.5), Cm(3.5), Cm(1.4), "6006\nEötvös", MID_GRAY, 9)
    draw_node(slide, Cm(21), Cm(10.5), Cm(3.5), Cm(1.4), "6007\n电梯实验", MID_GRAY, 9)

    draw_node(
        slide, Cm(17), Cm(13), Cm(4.5), Cm(1.6), "6008 等效原理\nbelief≈0.85", ACCENT_BLUE, 10
    )

    draw_node(slide, Cm(22.5), Cm(13), Cm(4), Cm(1.6), "6009\n光必须弯曲", ACCENT_BLUE, 10)

    draw_node(
        slide,
        Cm(20.5),
        Cm(15.5),
        Cm(5),
        Cm(1.8),
        '6010 Einstein\n0.87" belief≈0.8',
        ACCENT_BLUE,
        11,
    )

    draw_arrow(slide, Cm(18.75), Cm(11.9), Cm(19.25), Cm(13), MID_GRAY)
    draw_arrow(slide, Cm(22.75), Cm(11.9), Cm(20), Cm(13), MID_GRAY)
    draw_arrow(slide, Cm(19.25), Cm(14.6), Cm(24), Cm(13.6), MID_GRAY)
    draw_arrow(slide, Cm(24.5), Cm(14.6), Cm(23), Cm(15.5), MID_GRAY)

    # Highlight: same prediction, no contradiction
    from pptx.enum.shapes import MSO_SHAPE

    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(27.5), Cm(9), Cm(5.5), Cm(4.5))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xFD, 0xF0, 0xE0)
    box.line.color.rgb = ORANGE
    box.line.width = Pt(1.5)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Cm(0.3)
    tf.margin_top = Cm(0.3)
    p = tf.paragraphs[0]
    p.text = '两条路径\n→ 相同预测\n0.87"\n\n无 Contradiction\nEdge！'
    p.font.size = Pt(13)
    p.font.color.rgb = ORANGE
    p.font.bold = True
    p.font.name = FONT_BODY
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE


def slide_16_einstein_pkg4(prs):
    """Slide 16: Einstein Pkg 4 — GR divergence."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_header(slide, "Part 4 · 爱因斯坦思想实验", "Pkg 4: 广义相对论 — 定量矛盾出现")

    # ── Left: command box ──
    add_command_box(
        slide,
        Cm(1.5),
        Cm(4.2),
        Cm(14),
        Cm(14),
        [
            ("# Pkg 4: einstein1915_general_relativity", CODE_COMMENT),
            ('AddNode 6012 "GR: G_μν=8πGT_μν"  prior=0.85', CODE_KEYWORD),
            ('AddNode 6013 "光沿零测地线"        prior=0.85', CODE_DEFAULT),
            ("  (时间弯曲 + 空间弯曲 = 2个分量)", CODE_COMMENT),
            ('AddNode 6014 "GR预测: 1.75″"      prior=0.85', CODE_KEYWORD),
            ("  (= 2 × 0.87″, 空间弯曲贡献一半)", CODE_COMMENT),
            ('AddNode 6015 "水星进动 43″/世纪"    prior=0.9', CODE_DEFAULT),
            ("", CODE_DEFAULT),
            ("AddEdge 6006 deduction", CODE_DEFAULT),
            ("  [6012,6008] → [6013]  P=0.85", CODE_DEFAULT),
            ("AddEdge 6007 deduction", CODE_DEFAULT),
            ("  [6013] → [6014]       P=0.85", CODE_DEFAULT),
            ("", CODE_DEFAULT),
            ("AddEdge 6008 contradiction", CODE_RED),
            ("  [6014: 1.75″] ↔ [6004: 0.87″]", CODE_RED),
            ("  # 精确 2× 差异!", CODE_RED),
            ("", CODE_DEFAULT),
            ("AddEdge 6009 deduction", CODE_DEFAULT),
            ("  [6012] → [6015]       P=0.9", CODE_DEFAULT),
            ("", CODE_DEFAULT),
            ("# BP: 6014→0.85↑  6004→0.40↓  矛盾传播!", CODE_WARN),
        ],
    )

    # ── Right: graph ──
    # GR framework
    draw_node(
        slide, Cm(17.5), Cm(4.5), Cm(5), Cm(1.8), "6012 广义相对论\nbelief=0.85", ACCENT_BLUE, 11
    )
    draw_node(slide, Cm(24), Cm(4.5), Cm(4.5), Cm(1.8), "6008 等效原理\nbelief=0.85", MID_GRAY, 10)

    # Geodesics
    draw_node(
        slide,
        Cm(19),
        Cm(7.5),
        Cm(6),
        Cm(1.8),
        "6013 零测地线 (时间+空间)\nbelief=0.85",
        ACCENT_BLUE,
        10,
    )
    draw_arrow(slide, Cm(20), Cm(6.3), Cm(21), Cm(7.5), MID_GRAY)
    draw_arrow(slide, Cm(26.25), Cm(6.3), Cm(23), Cm(7.5), MID_GRAY)

    # GR prediction
    draw_node(slide, Cm(17), Cm(10.5), Cm(5.5), Cm(2.2), '6014\nGR: 1.75"\nbelief 0.85', GREEN, 13)
    draw_arrow(slide, Cm(22), Cm(9.3), Cm(20), Cm(10.5), MID_GRAY)

    # Newton prediction
    draw_node(
        slide,
        Cm(25.5),
        Cm(10.5),
        Cm(5.5),
        Cm(2.2),
        '6004\nNewton: 0.87"\nbelief 0.60→0.40↓',
        RED,
        12,
        WHITE,
    )

    # CONTRADICTION between 6014 and 6004
    draw_contradiction(slide, Cm(22.5), Cm(11.6), Cm(25.5), Cm(11.6))

    from pptx.enum.shapes import MSO_SHAPE

    cbox = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(22.8), Cm(13), Cm(5.5), Cm(1.8))
    cbox.fill.solid()
    cbox.fill.fore_color.rgb = RGBColor(0xFD, 0xE0, 0xE0)
    cbox.line.color.rgb = RED
    cbox.line.width = Pt(1)
    tf = cbox.text_frame
    p = tf.paragraphs[0]
    p.text = 'CONTRADICTION\n1.75" vs 0.87" (2×)'
    p.font.size = Pt(11)
    p.font.color.rgb = RED
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Mercury bonus
    draw_node(slide, Cm(17.5), Cm(15), Cm(5), Cm(1.8), '6015 水星进动\n43"/century  ✓', GREEN, 10)
    draw_arrow(slide, Cm(20), Cm(6.3), Cm(20), Cm(15), LIGHT_GRAY)

    add_text_box(
        slide,
        Cm(24),
        Cm(15.5),
        Cm(9),
        Cm(1.5),
        "GR 两个独立预测成功 → 6012 belief 强化",
        font_size=12,
        color=MID_GRAY,
    )


def slide_17_einstein_pkg5(prs):
    """Slide 17: Einstein Pkg 5 — Eddington observation + final beliefs."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_header(
        slide, "Part 4 · 爱因斯坦思想实验", "Pkg 5: Eddington 1919 — 判决 + Belief 全景"
    )

    # ── Left: command box ──
    add_command_box(
        slide,
        Cm(1.5),
        Cm(4.2),
        Cm(14),
        Cm(8),
        [
            ("# Pkg 5: eddington1919_solar_eclipse", CODE_COMMENT),
            ('AddNode 6017 "日食远征(两支队伍)"  prior=0.95', CODE_DEFAULT),
            ('AddNode 6018 "观测 1.61±0.30″ /  prior=0.9', CODE_DEFAULT),
            ('         1.98±0.16″"', CODE_DEFAULT),
            ('AddNode 6019 "观测支持GR排除Newton" prior=0.9', CODE_KEYWORD),
            ('AddNode 6020 "GR确认为优越理论"    prior=0.9', CODE_KEYWORD),
            ('AddNode 6021 "Newton降级为近似"    prior=0.9', CODE_DEFAULT),
            ("", CODE_DEFAULT),
            ("AddEdge 6011 deduction", CODE_DEFAULT),
            ("  [6017,6018] → [6019]  P=0.9", CODE_DEFAULT),
            ("AddEdge 6014 contradiction", CODE_RED),
            ("  [6019: ~1.7″] ↔ [6004: 0.87″]", CODE_RED),
            ("", CODE_DEFAULT),
            ("# BP: 6004→0.10↓↓  6014→0.95↑  6002→0.80(降级)", CODE_WARN),
        ],
    )

    # ── Left: final belief table ──
    rows = [
        ["节点", "含义", "初始", "最终"],
        ["6004", 'Soldner 0.87"', "0.60", "0.10 ↓↓"],
        ["6014", 'GR 1.75"', "0.85", "0.95 ↑↑"],
        ["6012", "广义相对论", "0.85", "0.95 ↑"],
        ["6002", "Newton 引力", "0.95", "0.80 (降级)"],
        ["6020", "GR confirmed", "0.90", "0.95 ↑"],
    ]
    make_table(
        slide,
        Cm(1.5),
        Cm(13),
        Cm(14),
        rows,
        col_widths=[Cm(2.5), Cm(5), Cm(2.5), Cm(4)],
        font_size=12,
    )

    # ── Right: final graph ──
    add_text_box(
        slide,
        Cm(17),
        Cm(4.2),
        Cm(10),
        Cm(1.2),
        "观测判决后的最终图:",
        font_size=16,
        color=BLACK,
        bold=True,
    )

    # Eddington observation
    draw_node(
        slide, Cm(18), Cm(5.5), Cm(6), Cm(1.8), '6018 观测: 1.61"/1.98"\n(≈1.7")', ACCENT_BLUE, 11
    )

    # GR prediction — confirmed
    draw_node(slide, Cm(17), Cm(9), Cm(5.5), Cm(2.2), '6014 GR: 1.75"\nbelief → 0.95 ↑↑', GREEN, 12)

    # Newton prediction — refuted
    draw_node(
        slide,
        Cm(25.5),
        Cm(9),
        Cm(5.5),
        Cm(2.2),
        '6004 Newton: 0.87"\nbelief → 0.10 ↓↓',
        RED,
        12,
        WHITE,
    )

    # Observation supports GR
    draw_arrow(slide, Cm(20), Cm(7.3), Cm(19.75), Cm(9), GREEN, Pt(2.5))
    add_text_box(
        slide, Cm(16.5), Cm(7.8), Cm(3), Cm(1), "✓ 吻合", font_size=11, color=GREEN, bold=True
    )

    # Observation contradicts Newton
    draw_contradiction(slide, Cm(24), Cm(7.3), Cm(28.25), Cm(9))
    add_text_box(
        slide, Cm(26), Cm(7.5), Cm(4), Cm(1), "✕ 2σ+ 偏离", font_size=11, color=RED, bold=True
    )

    # Contradiction edge between predictions
    draw_contradiction(slide, Cm(22.5), Cm(10.1), Cm(25.5), Cm(10.1))

    # GR framework
    draw_node(slide, Cm(17), Cm(13), Cm(5.5), Cm(1.8), "6012 广义相对论\nbelief → 0.95", GREEN, 11)

    # Newton framework (demoted)
    draw_node(
        slide,
        Cm(25.5),
        Cm(13),
        Cm(5.5),
        Cm(1.8),
        "6002 Newton 引力\nbelief → 0.80 (降级)",
        ORANGE,
        10,
    )

    draw_arrow(slide, Cm(19.75), Cm(11.2), Cm(19.75), Cm(13), MID_GRAY)
    draw_arrow(slide, Cm(28.25), Cm(11.2), Cm(28.25), Cm(13), MID_GRAY)

    # Key insight
    add_text_box(
        slide,
        Cm(17),
        Cm(15.5),
        Cm(16),
        Cm(2),
        "理论演替而非覆盖: Newton 降级为弱场近似，不删除",
        font_size=14,
        color=ACCENT_BLUE,
        bold=True,
    )


def slide_18_knowledge_package(prs):
    """Slide 18: Knowledge package = paper/theory."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_header(slide, "Part 5 · 知识包管理", "知识包 = 论文 / 理论")

    txBox = slide.shapes.add_textbox(Cm(2), Cm(4.5), Cm(16), Cm(10))
    tf = txBox.text_frame
    tf.word_wrap = True
    add_bullet_slide_text(
        tf,
        [
            "类比 Cargo / Julia Pkg：论文 = crate / package",
            "可版本化：同一知识可被新证据更新",
            "可声明依赖：推理链 = 依赖图",
            "可追踪变更：每次提交有完整历史",
        ],
        font_size=20,
    )

    # Package dependency diagram
    pkgs = [
        ("Aristotle\n350 BCE", Cm(20), Cm(5)),
        ("Galileo\n1638", Cm(20), Cm(8.5)),
        ("Medium\n1638", Cm(26), Cm(8.5)),
        ("Newton\n1687", Cm(23), Cm(12)),
        ("Apollo 15\n1971", Cm(29), Cm(12)),
    ]

    for name, x, y in pkgs:
        draw_node(slide, x, y, Cm(4.5), Cm(2), name, ACCENT_BLUE, 11)

    # Dependency arrows
    draw_arrow(slide, Cm(22.25), Cm(7), Cm(22.25), Cm(8.5), MID_GRAY)
    draw_arrow(slide, Cm(22.25), Cm(10.5), Cm(24.5), Cm(12), MID_GRAY)
    draw_arrow(slide, Cm(28.25), Cm(10.5), Cm(26), Cm(12), MID_GRAY)
    draw_arrow(slide, Cm(25.25), Cm(14), Cm(30), Cm(12), MID_GRAY)


def slide_19_isomorphism(prs):
    """Slide 19: Why package management and Gaia are isomorphic."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_header(slide, "Part 5 · 知识包管理", "为什么包管理和 Gaia 同构")

    rows = [
        ["", "Cargo / npm", "Gaia"],
        ["基本单元", "Package (crate)", "Knowledge Package (论文)"],
        ["依赖", "A depends on B", "edge tail → head"],
        ["版本", "semver", "Commit history"],
        ["冲突", "Version conflict", "Contradiction edge"],
        ["求解", "SAT solver / 拓扑排序", "Belief Propagation"],
        ["共同基础", "Horn clause 上的依赖传播", "Horn clause 上的依赖传播"],
    ]

    make_table(
        slide, Cm(2), Cm(4.5), Cm(30), rows, col_widths=[Cm(4.5), Cm(12), Cm(13.5)], font_size=15
    )

    add_text_box(
        slide,
        Cm(2),
        Cm(14.5),
        Cm(30),
        Cm(2),
        "包层面用 Cargo 拓扑算法 · 推理层面用 BP 算法 · 底层结构同构",
        font_size=18,
        color=ACCENT_BLUE,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )


def slide_20_environment(prs):
    """Slide 20: Environment = lightweight branch."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_header(slide, "Part 5 · 知识包管理", "Environment = 轻量分支")

    txBox = slide.shapes.add_textbox(Cm(2), Cm(4.5), Cm(16), Cm(10))
    tf = txBox.text_frame
    tf.word_wrap = True
    add_bullet_slide_text(
        tf,
        [
            "创建分支做思想实验（base snapshot + sparse overlay）",
            "伽利略例子 = 6 个 package 按序提交",
            "爱因斯坦例子 = 5 个 package 按序提交",
            "合并回主分支 = merge commit",
        ],
        font_size=20,
    )

    # Branch diagram
    from pptx.enum.shapes import MSO_SHAPE

    # Main branch line
    line = slide.shapes.add_connector(1, Cm(19), Cm(7), Cm(32), Cm(7))
    line.line.color.rgb = ACCENT_BLUE
    line.line.width = Pt(3)

    add_text_box(
        slide, Cm(18), Cm(5.5), Cm(3), Cm(1), "main", font_size=14, color=ACCENT_BLUE, bold=True
    )

    # Branch off
    line2 = slide.shapes.add_connector(1, Cm(21), Cm(7), Cm(23), Cm(10))
    line2.line.color.rgb = GREEN
    line2.line.width = Pt(2)
    line3 = slide.shapes.add_connector(1, Cm(23), Cm(10), Cm(30), Cm(10))
    line3.line.color.rgb = GREEN
    line3.line.width = Pt(2)

    add_text_box(
        slide,
        Cm(23),
        Cm(10.5),
        Cm(8),
        Cm(1),
        "experiment branch",
        font_size=13,
        color=GREEN,
        bold=True,
    )

    # Merge back
    line4 = slide.shapes.add_connector(1, Cm(30), Cm(10), Cm(32), Cm(7))
    line4.line.color.rgb = GREEN
    line4.line.width = Pt(2)

    # Commits on branch
    for i, label in enumerate(["Pkg 1", "Pkg 2", "...", "Pkg N"]):
        x = Cm(23.5 + i * 2)
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, Cm(9.6), Cm(0.8), Cm(0.8))
        dot.fill.solid()
        dot.fill.fore_color.rgb = GREEN
        dot.line.fill.background()

    # Commits on main
    for i in range(4):
        x = Cm(19.5 + i * 3.5)
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, Cm(6.6), Cm(0.8), Cm(0.8))
        dot.fill.solid()
        dot.fill.fore_color.rgb = ACCENT_BLUE
        dot.line.fill.background()


def slide_21_lkm_vision(prs):
    """Slide 21: Large Knowledge Model vision."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_header(slide, "Part 6 · 远景与总结", "Large Knowledge Model — 远景")

    rows = [
        ["", "Large Language Model", "Large Knowledge Model"],
        ["存储形式", "压缩在参数中", "显式图结构 (每条可寻址)"],
        ["可追溯性", "无法追溯来源", "每条知识有完整推理链"],
        ["可验证性", "黑箱", "每步推理可审计"],
        ["矛盾处理", "幻觉 (Hallucination)", "显式 Contradiction Edge"],
        ["更新方式", "重新训练", "增量提交 + BP 传播"],
        ["规模", "~10¹² 参数", "10⁹ 节点, 5×10⁹ 超边"],
    ]

    make_table(
        slide, Cm(2), Cm(4.5), Cm(30), rows, col_widths=[Cm(4.5), Cm(12), Cm(13.5)], font_size=15
    )


def slide_22_scale(prs):
    """Slide 22: How to reach billion scale."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_header(slide, "Part 6 · 远景与总结", "十亿级规模怎么做")

    # Three storage layers
    from pptx.enum.shapes import MSO_SHAPE

    layers = [
        ("LanceDB", "内容 + 元数据 + BM25 全文", ACCENT_BLUE),
        ("Neo4j", "图拓扑 + 超边关系", ORANGE),
        ("Vector Index", "语义嵌入相似搜索", GREEN),
    ]

    for i, (name, desc, color) in enumerate(layers):
        x = Cm(2 + i * 10.5)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Cm(4.5), Cm(9), Cm(3))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"{name}\n{desc}"
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    add_text_box(
        slide,
        Cm(2),
        Cm(8),
        Cm(30),
        Cm(1),
        "三层存储：每一层负责一种查询模式，互为补充",
        font_size=16,
        color=MID_GRAY,
        alignment=PP_ALIGN.CENTER,
    )

    txBox = slide.shapes.add_textbox(Cm(2), Cm(9.5), Cm(15), Cm(8))
    tf = txBox.text_frame
    tf.word_wrap = True
    add_bullet_slide_text(
        tf,
        [
            "分布式 BP: 图分区 + 残差调度 + 层次化",
            "预接地 (Pre-grounding): 避免 MLN 组合爆炸",
            "增量 BP: 只传播受影响的子图",
            "近似推断: Loopy BP 收敛性保证",
        ],
        font_size=20,
    )


def slide_23_usage(prs):
    """Slide 23: How agents use Gaia."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_header(slide, "Part 6 · 远景与总结", "大概怎么用")

    # Four usage scenarios
    from pptx.enum.shapes import MSO_SHAPE

    scenarios = [
        ("1. 查询知识", "向量 + BM25 + 拓扑\n三路召回 → 带 belief 的知识", ACCENT_BLUE),
        ("2. 提交发现", "Commit → Review → Merge\n→ BP 传播更新 belief", GREEN),
        ("3. 追溯推理", "任意结论 → 沿超边回溯\n→ 完整推理链", ORANGE),
        ("4. 发现矛盾", "Contradiction Edge 自动标记\n→ 影响传播到下游", RED),
    ]

    for i, (title, desc, color) in enumerate(scenarios):
        col = i % 2
        row = i // 2
        x = Cm(2 + col * 16)
        y = Cm(5 + row * 5.5)

        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Cm(14), Cm(4))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
        box.line.color.rgb = color
        box.line.width = Pt(2)
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Cm(0.5)
        tf.margin_top = Cm(0.3)
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(20)
        p.font.color.rgb = color
        p.font.bold = True
        p.font.name = FONT_BODY
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(15)
        p2.font.color.rgb = DARK_GRAY
        p2.font.name = FONT_BODY
        p2.space_before = Pt(6)


def slide_24_summary(prs):
    """Slide 24: Summary."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_header(slide, "Part 6 · 远景与总结", "总结")

    items = [
        "Gaia = 概率推理图 + 知识包管理 + Git 工作流",
        "矛盾是一等公民 — Contradiction Edge + BP 自动量化",
        "旧理论不删除，belief 下降 — 保留科学史完整性",
        "每条知识可寻址、可溯源、可验证",
        "Agent 的外部可审计记忆 — Agentic Science at Scale 的基础设施",
    ]

    txBox = slide.shapes.add_textbox(Cm(4), Cm(5), Cm(26), Cm(12))
    tf = txBox.text_frame
    tf.word_wrap = True
    add_bullet_slide_text(tf, items, font_size=24, spacing=Pt(16))


def slide_25_thankyou(prs):
    """Slide 25: Thank you."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    from pptx.enum.shapes import MSO_SHAPE

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Cm(0.6))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_BLUE
    bar.line.fill.background()

    add_text_box(
        slide,
        Cm(3),
        Cm(5),
        Cm(28),
        Cm(3),
        "Thank You",
        font_size=60,
        color=BLACK,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )

    add_text_box(
        slide,
        Cm(3),
        Cm(9),
        Cm(28),
        Cm(2),
        "Questions & Discussion",
        font_size=30,
        color=MID_GRAY,
        alignment=PP_ALIGN.CENTER,
    )

    add_text_box(
        slide,
        Cm(3),
        Cm(13),
        Cm(28),
        Cm(2),
        "github.com/SiliconEinstein/Gaia",
        font_size=18,
        color=ACCENT_BLUE,
        alignment=PP_ALIGN.CENTER,
    )


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Build all slides
    slide_01_cover(prs)
    slide_02_agent_science(prs)
    slide_03_bottleneck1(prs)
    slide_04_bottleneck2(prs)
    slide_05_core_idea(prs)
    slide_06_paradigm_shift(prs)
    slide_07_contradiction(prs)
    slide_08_bp(prs)
    slide_09_git_workflow(prs)
    slide_10_galileo_overview(prs)
    slide_11_galileo_pkg12(prs)
    slide_12_galileo_pkg345(prs)
    slide_13_galileo_pkg6(prs)
    slide_14_einstein_overview(prs)
    slide_15_einstein_pkg123(prs)
    slide_16_einstein_pkg4(prs)
    slide_17_einstein_pkg5(prs)
    slide_18_knowledge_package(prs)
    slide_19_isomorphism(prs)
    slide_20_environment(prs)
    slide_21_lkm_vision(prs)
    slide_22_scale(prs)
    slide_23_usage(prs)
    slide_24_summary(prs)
    slide_25_thankyou(prs)

    out_path = Path(__file__).resolve().parent.parent / "Gaia_Presentation.pptx"
    prs.save(str(out_path))
    print(f"Saved presentation to {out_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
