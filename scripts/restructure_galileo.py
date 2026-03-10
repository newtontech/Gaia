"""Restructure Galileo slides: 4 → 6 slides with correct logic and diagrams."""

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from copy import deepcopy
from lxml import etree

# ═══════════════════════════════════════════════════════════════
# Constants
# ═══════════════════════════════════════════════════════════════
NSMAP = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"

# Colors (hex strings)
WHITE = "E0E0E0"
CYAN = "5CC9F5"
GRAY = "808888"
GOLD = "FFD54F"
RED = "FF8A80"
GREEN = "69F0AE"
ORANGE = "FFAB40"
HEADER = "FFFFFF"

# Node fill colors
NODE_GRAY = RGBColor(0x66, 0x66, 0x66)      # neutral nodes
NODE_RED = RGBColor(0xD3, 0x2F, 0x2F)       # declining belief
NODE_BLUE = RGBColor(0x1B, 0x5E, 0x9F)      # deduction nodes
NODE_GREEN = RGBColor(0x2E, 0x7D, 0x32)     # conclusion / confirmed
NODE_ORANGE = RGBColor(0xE6, 0x51, 0x00)    # contradiction
NODE_GOLD = RGBColor(0xF9, 0xA8, 0x25)      # key conclusions (★)
LINE_GRAY = RGBColor(0x99, 0x99, 0x99)
LINE_RED = RGBColor(0xF4, 0x43, 0x36)       # contradiction line


# ═══════════════════════════════════════════════════════════════
# Colorization helpers (from colorize_slides.py)
# ═══════════════════════════════════════════════════════════════
def colorize_line(line):
    """Return list of (text, color_hex, bold) segments for a line."""
    stripped = line.strip()
    if not stripped:
        return [("", WHITE, False)]
    if "═══" in stripped:
        return [(line, HEADER, True)]
    if stripped.startswith("#") and "═══" not in stripped:
        if "CONTRADICTION" in stripped or "矛盾" in stripped:
            return [(line, ORANGE, False)]
        if "★" in stripped:
            return [(line, GOLD, True)]
        if "✓" in stripped:
            return [(line, GREEN, False)]
        return [(line, GRAY, False)]
    if stripped.startswith("$ gaia"):
        parts = []
        if "  #" in line:
            idx = line.index("  #")
            parts.append((line[:idx], CYAN, False))
            parts.append((line[idx:], GRAY, False))
        elif "# " in line and not line.strip().startswith("#"):
            idx = line.index("# ")
            parts.append((line[:idx], CYAN, False))
            parts.append((line[idx:], GRAY, False))
        else:
            parts.append((line, CYAN, False))
        return parts
    if "★" in line:
        return [(line, GOLD, True)]
    if "contradiction" in stripped.lower() or "CONTRADICTION" in stripped:
        return [(line, ORANGE, True)]
    if stripped.startswith("→") or stripped.startswith("→"):
        if "★" in line:
            return [(line, GOLD, True)]
        return [(line, GREEN, False)]
    has_down = "↓" in line
    has_up = "↑" in line
    if has_down and has_up:
        segments = []
        parts = line.split("  ")
        for part in parts:
            ps = part.strip()
            if not ps:
                continue
            if "↓" in part:
                segments.append(("  " + ps, RED, False))
            elif "↑" in part:
                segments.append(("  " + ps, GREEN, False))
            else:
                segments.append(("  " + ps, WHITE, False))
        return segments if segments else [(line, WHITE, False)]
    if has_down:
        return [(line, RED, False)]
    if has_up:
        return [(line, GREEN, False)]
    return [(line, WHITE, False)]


def make_run_xml(text, color_hex, bold=False, sz=None):
    """Create an <a:r> element with colored text."""
    ns = NS_A
    r = etree.Element(f"{{{ns}}}r")
    rPr = etree.SubElement(r, f"{{{ns}}}rPr")
    rPr.set("lang", "en-US")
    rPr.set("dirty", "0")
    if bold:
        rPr.set("b", "1")
    if sz:
        rPr.set("sz", str(sz))

    solidFill = etree.SubElement(rPr, f"{{{ns}}}solidFill")
    srgbClr = etree.SubElement(solidFill, f"{{{ns}}}srgbClr")
    srgbClr.set("val", color_hex)

    latin = etree.SubElement(rPr, f"{{{ns}}}latin")
    latin.set("typeface", "Helvetica Neue")

    t = etree.SubElement(r, f"{{{ns}}}t")
    t.text = text
    if text and (text[0] == " " or text[-1] == " "):
        t.set("{http://www.w3.org/XML/1998/namespace}space", "preserve")

    return r


def set_colored_text(text_frame, lines):
    """Replace text frame content with syntax-highlighted lines."""
    ref_p = text_frame.paragraphs[0]._p
    ref_pPr = ref_p.find("a:pPr", NSMAP)
    ref_pPr_xml = deepcopy(ref_pPr) if ref_pPr is not None else None

    txBody = text_frame._txBody
    for p in txBody.findall(f"{{{NS_A}}}p"):
        txBody.remove(p)

    for i, line in enumerate(lines):
        p = etree.SubElement(txBody, f"{{{NS_A}}}p")
        if ref_pPr_xml is not None:
            pPr = deepcopy(ref_pPr_xml)
            if i > 0 and "algn" in pPr.attrib:
                del pPr.attrib["algn"]
            p.insert(0, pPr)
        segments = colorize_line(line)
        for text, color, bold in segments:
            r = make_run_xml(text, color, bold)
            p.append(r)


# ═══════════════════════════════════════════════════════════════
# Slide manipulation helpers
# ═══════════════════════════════════════════════════════════════
def duplicate_slide(prs, source_idx):
    """Duplicate a slide to the end of the presentation. Returns new slide."""
    source = prs.slides[source_idx]
    new_slide = prs.slides.add_slide(source.slide_layout)

    # Clear auto-generated shapes
    sp_tree = new_slide.shapes._spTree
    for child in list(sp_tree)[2:]:  # keep nvGrpSpPr + grpSpPr
        sp_tree.remove(child)

    # Copy all shapes from source
    src_tree = source.shapes._spTree
    for child in list(src_tree)[2:]:
        sp_tree.append(deepcopy(child))

    # Copy background if present
    src_bg = source._element.find(f"{{{NS_P}}}bg")
    if src_bg is not None:
        new_bg = new_slide._element.find(f"{{{NS_P}}}bg")
        if new_bg is not None:
            new_slide._element.remove(new_bg)
        new_slide._element.insert(0, deepcopy(src_bg))

    return new_slide


def move_slide(prs, from_idx, to_idx):
    """Move slide from from_idx to before to_idx."""
    sldIdLst = prs.slides._sldIdLst
    slides = list(sldIdLst)
    el = slides[from_idx]
    sldIdLst.remove(el)
    slides = list(sldIdLst)
    if to_idx >= len(slides):
        sldIdLst.append(el)
    else:
        slides[to_idx].addprevious(el)


def clear_right_shapes(slide, x_threshold=5400000):
    """Remove all shapes whose left edge is at or beyond x_threshold."""
    sp_tree = slide.shapes._spTree
    to_remove = []
    for child in list(sp_tree)[2:]:
        # Get position
        spPr = child.find(f".//{{{NS_A}}}off")
        if spPr is not None:
            x = int(spPr.get("x", "0"))
            if x >= x_threshold:
                to_remove.append(child)
    for child in to_remove:
        sp_tree.remove(child)


# ═══════════════════════════════════════════════════════════════
# Diagram shape builders
# ═══════════════════════════════════════════════════════════════
def add_node_box(slide, left, top, width, height, lines, fill_color,
                 font_size=Pt(9), font_color=RGBColor(0xFF, 0xFF, 0xFF)):
    """Add a rounded rectangle node with multi-line text."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = fill_color
    shape.line.width = Emu(12700)

    # Set text
    tf = shape.text_frame
    tf.word_wrap = True
    for i, line_text in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line_text
        p.font.size = font_size
        p.font.color.rgb = font_color
        p.font.bold = (i == 0)  # First line bold
        p.space_after = Pt(0)
        p.space_before = Pt(0)
    # Center text vertically
    from pptx.enum.text import PP_ALIGN
    for p in tf.paragraphs:
        p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    return shape


def add_label(slide, left, top, width, height, text,
              font_size=Pt(9), font_color=RGBColor(0xFF, 0xFF, 0xFF),
              bold=False):
    """Add a text label (no fill)."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.space_after = Pt(0)
    from pptx.enum.text import PP_ALIGN
    p.alignment = PP_ALIGN.CENTER
    return txBox


def add_line(slide, x1, y1, x2, y2, color=LINE_GRAY, width=Emu(12700), dashed=False):
    """Add a connector line."""
    connector = slide.shapes.add_connector(
        1,  # cxnSp type (straight)
        x1, y1, x2, y2
    )
    connector.line.color.rgb = color
    connector.line.width = width
    if dashed:
        # Set dash style via XML
        ln = connector._element.find(f".//{{{NS_A}}}ln")
        if ln is not None:
            prstDash = etree.SubElement(ln, f"{{{NS_A}}}prstDash")
            prstDash.set("val", "dash")
    return connector


def add_arrow(slide, x1, y1, x2, y2, color=LINE_GRAY, width=Emu(12700)):
    """Add an arrow (line with arrowhead)."""
    connector = add_line(slide, x1, y1, x2, y2, color, width)
    # Add arrowhead
    ln = connector._element.find(f".//{{{NS_A}}}ln")
    if ln is None:
        spPr = connector._element.find(f".//{{{NS_A}}}spPr")
        ln = etree.SubElement(spPr, f"{{{NS_A}}}ln")
    tailEnd = etree.SubElement(ln, f"{{{NS_A}}}tailEnd")
    tailEnd.set("type", "triangle")
    tailEnd.set("w", "med")
    tailEnd.set("len", "med")
    return connector


# ═══════════════════════════════════════════════════════════════
# Slide content definitions
# ═══════════════════════════════════════════════════════════════

# Slide 10 (idx 9): Overview — keep mostly as-is, update text
OVERVIEW_TEXT = [
    "1638 年伽利略《Discorsi》: 一个思想实验推翻 2000 年学说",
    "",
    "6 个 knowledge package，跨越 2300 年:",
    "  Pkg 1  亚里士多德先验 → 建立 v∝W (0.70)",
    "  Pkg 2  绑球悖论 → Contradiction① → v∝W 降至 0.35",
    "  Pkg 3  密度观察 → 空气阻力才是真因",
    "  Pkg 4  真空预测 → 去掉空气，所有物体等速",
    "  Pkg 5  牛顿 F=ma → Contradiction② → v∝W 降至 0.12",
    "  Pkg 6  Apollo 15 → 月球真空实验 → v∝W 降至 0.05",
    "",
    "演示 Gaia 核心机制:",
    "  • contradiction edge 驱动 belief 修正",
    "  • 三条独立证据线汇聚 → 高 belief",
    "  • 旧理论不删除，belief 自然下降",
]

# Slide 11 (Pkg 1): Aristotle
PKG1_CODE = [
    "$ gaia init galileo_tied_balls",
    "",
    "# ═══ Pkg 1: aristotle_physics ═══",
    "# 前提: 亚里士多德自然运动学说 (c.350 BCE)",
    "",
    '$ gaia node add "自然运动学说: 重物自然趋向下方,',
    '    下落速度取决于重量" \\',
    "    --prior 0.90 --type paper-extract  → 5001",
    "",
    '$ gaia node add "日常观察: 石头比树叶先落地" \\',
    "    --prior 0.95                       → 5002",
    "",
    "# 从学说+观察 归纳出一般定律",
    "$ gaia edge add --tail 5001,5002 \\",
    "    --head 5003 --type abstraction",
    '  → 5003 "v ∝ W: 越重越快" prior=0.70',
    "",
    "$ gaia commit -m \"aristotle_physics\"",
    "$ gaia propagate",
    "  5003 (v ∝ W): belief = 0.70",
    "# ✓ 这就是 2000 年间公认的物理学常识",
]

# Slide 12 (Pkg 2): Tied Balls — THE KEY
PKG2_CODE = [
    "# ═══ Pkg 2: galileo1638_tied_balls ═══",
    "# 核心思想实验: 绑球悖论",
    "",
    '$ gaia node add "设定: 重球H绑轻球L" \\',
    "    --prior 0.99                       → 5004",
    "",
    "# 推导A: L比H慢，L拖慢H → 组合HL比H更慢",
    "$ gaia edge add --tail 5003,5004 \\",
    "    --head 5005 --type deduction",
    '  → 5005 "HL比H更慢" (L拖慢了H)',
    "",
    "# 推导B: HL总重>H → 组合HL比H更快",
    "$ gaia edge add --tail 5003,5004 \\",
    "    --head 5006 --type deduction",
    '  → 5006 "HL比H更快" (更重=更快)',
    "",
    "# ★ 同一物体不可能既更快又更慢!",
    "$ gaia edge add --tail 5005,5006 \\",
    "    --type contradiction",
    "",
    "$ gaia commit && gaia propagate",
    "  5003 (v∝W): 0.70→0.35 ↓  矛盾回传到前提!",
    "  5008 (v∝W自相矛盾): belief=0.82 ↑",
]

# Slide 13 (Pkg 3-4): Medium + Vacuum
PKG34_CODE = [
    "# ═══ Pkg 3: medium_density ═══",
    "# 伽利略的下一步: 解释为何亚里士多德\"看起来\"对",
    "",
    '$ gaia node add "观察: 水中落速差异 > 空气中" \\',
    "    --prior 0.90                       → 5009",
    '$ gaia node add "规律: 介质越稀，差异越小" \\',
    "    --prior 0.85                       → 5010",
    "$ gaia edge add --tail 5009,5010 \\",
    "    --head 5011 --type deduction",
    "  → 5011 \"★ 空气阻力才是混淆因素\"",
    "  # 石头快于叶子不是因为更重，而是空气阻力不同!",
    "",
    "# ═══ Pkg 4: vacuum_prediction ═══",
    "# 逻辑推论: 若差异来自介质 → 去掉介质则差异消失",
    "",
    "$ gaia edge add --tail 5008,5011 \\",
    "    --head 5012 --type deduction",
    "  → 5012 \"★ 真空中所有物体等速下落\"",
    '$ gaia node add "斜面实验: 不同球同时到底" \\',
    "    --prior 0.90  # 部分验证            → 5013",
    "",
    "$ gaia commit && gaia propagate",
    "  5012 (真空等速): belief=0.78 ↑",
    "  5003 (v∝W): 0.35→0.28 ↓  替代解释进一步削弱",
]

# Slide 14 (Pkg 5): Newton
PKG5_CODE = [
    "# ═══ Pkg 5: newton1687_principia ═══",
    "# 50年后: 牛顿从第一性原理独立推导",
    "",
    '$ gaia node add "F = ma (牛顿第二定律)" \\',
    "    --prior 0.95                       → 5015",
    '$ gaia node add "F = mg (万有引力)" \\',
    "    --prior 0.95                       → 5016",
    "",
    "# F=ma=mg → 两边除以m → a=g",
    "$ gaia edge add --tail 5015,5016 \\",
    "    --head 5017 --type deduction",
    '  → 5017 "a = g: 加速度与质量无关!"',
    "",
    "# ★ 第二条矛盾线: Newton a=g vs Aristotle v∝W",
    "$ gaia edge add --tail 5003,5017 \\",
    "    --type contradiction",
    "",
    "# 牛顿独立确认伽利略的真空预测",
    "$ gaia edge add --tail 5017 \\",
    "    --head 5012 --type deduction",
    "",
    "$ gaia commit && gaia propagate",
    "  5003 (v∝W): 0.28→0.12 ↓↓  两条矛盾线夹击",
    "  5017 (a=g): belief=0.93 ↑",
    "  5012 (真空等速): 0.78→0.87 ↑  理论+逻辑双支撑",
]

# Slide 15 (Pkg 6 + Summary): Apollo
PKG6_CODE = [
    "# ═══ Pkg 6: apollo15_feather_drop ═══",
    "# 1971: 人类首次在真空中直接验证",
    "",
    '$ gaia node add "月球: 近乎真空" \\',
    "    --prior 0.99                       → 5018",
    '$ gaia node add "锤子(1.32kg) = 羽毛(0.03kg)',
    '    同时落地 (质量比44:1!)" \\',
    "    --prior 0.99                       → 5019",
    "",
    "$ gaia edge add --tail 5018,5019 \\",
    "    --head 5020 --type deduction",
    "",
    "$ gaia commit -m \"apollo15\"",
    "$ gaia propagate",
    "  5003 (v∝W): 0.12→0.05 ↓  几乎归零",
    "  5012 (真空等速): 0.87→0.95 ↑  三线汇聚",
    "  5017 (a=g): 0.93→0.96 ↑  理论+实验",
    "  5020 (月球验证): belief=0.98 ↑  决定性",
    "",
    "$ gaia test     # 验证 beliefs.yaml",
    "$ gaia publish  # → 远程 registry",
]


# ═══════════════════════════════════════════════════════════════
# Diagram builders for each slide
# ═══════════════════════════════════════════════════════════════

# Common dimensions (EMU)
E = Emu  # alias

# Right panel: x ∈ [5760000, 11520000], y ∈ [1500000, 6200000]
RX = 5760000  # right panel left edge
RCX = 8640000  # center x
RW = 5760000  # total width


def build_diagram_pkg1(slide):
    """Pkg 1: Aristotle — 5001 + 5002 → 5003."""
    clear_right_shapes(slide)

    # Title
    add_label(slide, RX, E(1500000), E(RW), E(360000),
              "亚里士多德先验知识图", Pt(14),
              RGBColor(0xFF, 0xFF, 0xFF), bold=True)

    # Node 5001
    add_node_box(slide, E(5940000), E(2160000), E(2340000), E(720000),
                      ["5001", "自然运动学说: 重物趋下"],
                      NODE_GRAY, Pt(9))

    # Node 5002
    add_node_box(slide, E(8820000), E(2160000), E(2340000), E(720000),
                      ["5002", "日常观察: 石头>叶子"],
                      NODE_GRAY, Pt(9))

    # Arrow 5001 → 5003
    add_arrow(slide, E(7110000), E(2880000), E(7920000), E(3600000))
    # Arrow 5002 → 5003
    add_arrow(slide, E(9990000), E(2880000), E(9360000), E(3600000))

    # Edge label
    add_label(slide, E(8100000), E(3060000), E(1800000), E(360000),
              "abstraction", Pt(8), RGBColor(0x99, 0x99, 0x99))

    # Node 5003 — the key node
    add_node_box(slide, E(7200000), E(3600000), E(2880000), E(900000),
                      ["5003: v ∝ W 定律", "越重越快", "belief = 0.70"],
                      NODE_RED, Pt(10))

    # Summary text at bottom
    add_label(slide, RX, E(5040000), E(RW), E(720000),
              "这是 2000 年间被广泛接受的\"常识\"\n"
              "接下来伽利略将用纯推理推翻它",
              Pt(10), RGBColor(0xBB, 0xBB, 0xBB))


def build_diagram_pkg2(slide):
    """Pkg 2: Tied Balls — the key contradiction."""
    clear_right_shapes(slide)

    # Title
    add_label(slide, RX, E(1500000), E(RW), E(360000),
              "绑球悖论: 同一前提推出矛盾结论", Pt(13),
              RGBColor(0xFF, 0xFF, 0xFF), bold=True)

    # Top row: 5003 (v∝W) + 5004 (绑球设定)
    add_node_box(slide, E(5940000), E(2016000), E(1980000), E(648000),
                 ["5003: v∝W", "belief=0.70"],
                 NODE_RED, Pt(9))

    add_node_box(slide, E(9180000), E(2016000), E(1980000), E(648000),
                 ["5004: 设定", "重球H绑轻球L"],
                 NODE_GRAY, Pt(9))

    # Arrows down to deductions
    add_arrow(slide, E(6930000), E(2664000), E(6570000), E(3240000))  # →5005
    add_arrow(slide, E(6930000), E(2664000), E(9810000), E(3240000))  # →5006
    add_arrow(slide, E(10170000), E(2664000), E(6930000), E(3240000))  # →5005
    add_arrow(slide, E(10170000), E(2664000), E(10170000), E(3240000))  # →5006

    # Deduction A: 5005
    add_node_box(slide, E(5580000), E(3240000), E(2340000), E(756000),
                 ["5005: 推导A", "L拖慢H → HL更慢"],
                 NODE_BLUE, Pt(9))

    # Deduction B: 5006
    add_node_box(slide, E(9000000), E(3240000), E(2340000), E(756000),
                 ["5006: 推导B", "HL更重 → HL更快"],
                 NODE_BLUE, Pt(9))

    # Contradiction line between 5005 and 5006
    add_line(slide, E(7920000), E(3618000), E(9000000), E(3618000),
             LINE_RED, Emu(25400), dashed=True)

    # Contradiction label
    add_label(slide, E(7740000), E(3996000), E(1800000), E(432000),
              "CONTRADICTION ①\n同一物体不可能\n既更快又更慢!",
              Pt(9), RGBColor(0xF4, 0x43, 0x36), bold=True)

    # Arrow down to conclusion
    add_arrow(slide, E(8280000), E(4428000), E(8280000), E(4860000))

    # Node 5008: conclusion
    add_node_box(slide, E(6840000), E(4860000), E(3240000), E(756000),
                 ["5008: v∝W 必然错误!", "两条推导都合法，前提必须有错"],
                 NODE_GREEN, Pt(9))

    # Result annotation
    add_label(slide, RX, E(5760000), E(RW), E(432000),
              "5003 (v∝W): 0.70 → 0.35 ↓   矛盾回传到共享前提",
              Pt(10), RGBColor(0xFF, 0x8A, 0x80))


def build_diagram_pkg34(slide):
    """Pkg 3-4: Medium density + Vacuum prediction."""
    clear_right_shapes(slide)

    # Title
    add_label(slide, RX, E(1500000), E(RW), E(360000),
              "找到真正原因 → 做出可验证的预测", Pt(13),
              RGBColor(0xFF, 0xFF, 0xFF), bold=True)

    # --- Pkg 3 section ---
    add_label(slide, E(5940000), E(1908000), E(1800000), E(288000),
              "Pkg 3: 密度观察", Pt(10),
              RGBColor(0x99, 0x99, 0x99), bold=True)

    # Nodes 5009 + 5010
    add_node_box(slide, E(5940000), E(2196000), E(2160000), E(576000),
                 ["5009: 水中差异更大", "5010: 介质越稀差异越小"],
                 NODE_GRAY, Pt(8))

    # Arrow to 5011
    add_arrow(slide, E(7020000), E(2772000), E(7020000), E(3132000))

    # Node 5011 — ★ key conclusion 1
    add_node_box(slide, E(5940000), E(3132000), E(2520000), E(648000),
                 ["★ 5011: 空气阻力", "才是混淆因素"],
                 NODE_GOLD, Pt(9))

    # --- Pkg 4 section ---
    add_label(slide, E(8820000), E(1908000), E(2520000), E(288000),
              "Pkg 4: 真空预测", Pt(10),
              RGBColor(0x99, 0x99, 0x99), bold=True)

    # Node 5008 (from Pkg 2)
    add_node_box(slide, E(9000000), E(2196000), E(2160000), E(576000),
                 ["5008: v∝W 错误", "(来自 Pkg 2)"],
                 NODE_GREEN, Pt(8))

    # Arrows to 5012
    add_arrow(slide, E(7560000), E(3456000), E(8280000), E(4032000))
    add_arrow(slide, E(10080000), E(2772000), E(9360000), E(4032000))

    # Edge label
    add_label(slide, E(8280000), E(3528000), E(1800000), E(360000),
              "v∝W错 + 空气阻力\n→ 去掉空气 = 等速",
              Pt(8), RGBColor(0x99, 0x99, 0x99))

    # Node 5012 — ★ key conclusion 2
    add_node_box(slide, E(7560000), E(4032000), E(3060000), E(648000),
                 ["★ 5012: 真空等速预测", "真空中所有物体等速下落"],
                 NODE_GOLD, Pt(9))

    # Node 5013 — inclined plane
    add_node_box(slide, E(7560000), E(4968000), E(2340000), E(576000),
                 ["5013: 斜面实验", "部分验证 ✓"],
                 NODE_GRAY, Pt(8))

    add_arrow(slide, E(8730000), E(4680000), E(8730000), E(4968000))

    # Result annotation
    add_label(slide, RX, E(5688000), E(RW), E(504000),
              "伽利略两大结论:\n"
              "① 石头快于叶子是因为空气阻力，不是因为更重\n"
              "② 在真空中，所有物体等速下落",
              Pt(10), RGBColor(0xFF, 0xD5, 0x4F))


def build_diagram_pkg5(slide):
    """Pkg 5: Newton — second contradiction line."""
    clear_right_shapes(slide)

    # Title
    add_label(slide, RX, E(1500000), E(RW), E(360000),
              "独立的理论推导: 第二条矛盾线", Pt(13),
              RGBColor(0xFF, 0xFF, 0xFF), bold=True)

    # F=ma
    add_node_box(slide, E(6120000), E(2088000), E(1800000), E(576000),
                 ["5015: F = ma", "牛顿第二定律"],
                 NODE_GRAY, Pt(9))

    # F=mg
    add_node_box(slide, E(9180000), E(2088000), E(1800000), E(576000),
                 ["5016: F = mg", "万有引力"],
                 NODE_GRAY, Pt(9))

    # Arrows to a=g
    add_arrow(slide, E(7020000), E(2664000), E(7920000), E(3168000))
    add_arrow(slide, E(10080000), E(2664000), E(9360000), E(3168000))

    # Label: cancellation
    add_label(slide, E(7740000), E(2700000), E(1800000), E(360000),
              "ma = mg → a = g\n质量约掉了!",
              Pt(9), RGBColor(0x69, 0xF0, 0xAE))

    # Node 5017: a=g
    add_node_box(slide, E(7200000), E(3168000), E(2880000), E(720000),
                 ["5017: a = g", "加速度与质量无关!"],
                 NODE_GREEN, Pt(10))

    # Contradiction with 5003
    add_node_box(slide, E(5940000), E(4392000), E(1980000), E(648000),
                 ["5003: v∝W", "0.28 → 0.12 ↓↓"],
                 NODE_RED, Pt(9))

    add_line(slide, E(7920000), E(4716000), E(7920000), E(3888000),
             LINE_RED, Emu(25400), dashed=True)

    add_label(slide, E(7740000), E(4032000), E(2160000), E(432000),
              "CONTRADICTION ②\nNewton a=g vs\nAristotle v∝W",
              Pt(9), RGBColor(0xF4, 0x43, 0x36), bold=True)

    # Arrow to 5012 (confirm vacuum prediction)
    add_node_box(slide, E(9180000), E(4392000), E(2160000), E(648000),
                 ["5012: 真空等速", "0.78 → 0.87 ↑"],
                 NODE_GOLD, Pt(9))

    add_arrow(slide, E(9360000), E(3888000), E(10260000), E(4392000))

    add_label(slide, E(9360000), E(3888000), E(1800000), E(360000),
              "独立理论确认",
              Pt(8), RGBColor(0x99, 0x99, 0x99))

    # Summary
    add_label(slide, RX, E(5400000), E(RW), E(720000),
              "牛顿力学从 F=ma, F=mg 独立推导出 a=g\n"
              "这是第二条完全独立的矛盾线\n"
              "同时从理论上确认了伽利略的真空等速预测",
              Pt(10), RGBColor(0xBB, 0xBB, 0xBB))


def build_diagram_pkg6(slide):
    """Pkg 6: Apollo 15 + convergence summary."""
    clear_right_shapes(slide)

    # Title
    add_label(slide, RX, E(1500000), E(RW), E(360000),
              "三条独立证据线汇聚", Pt(14),
              RGBColor(0xFF, 0xFF, 0xFF), bold=True)

    # Evidence line 1: Logical
    add_node_box(slide, E(6120000), E(2088000), E(1620000), E(648000),
                 ["① 逻辑矛盾", "绑球悖论"],
                 NODE_ORANGE, Pt(9))
    add_label(slide, E(7920000), E(2088000), E(2700000), E(648000),
              "Pkg 2: 纯推理\n从v∝W自身推出矛盾\nContradiction Edge ①",
              Pt(8), RGBColor(0xBB, 0xBB, 0xBB))

    # Evidence line 2: Theoretical
    add_node_box(slide, E(6120000), E(3060000), E(1620000), E(648000),
                 ["② 理论推导", "Newton a=g"],
                 NODE_ORANGE, Pt(9))
    add_label(slide, E(7920000), E(3060000), E(2700000), E(648000),
              "Pkg 5: F=ma, F=mg → a=g\n从第一性原理独立推导\nContradiction Edge ②",
              Pt(8), RGBColor(0xBB, 0xBB, 0xBB))

    # Evidence line 3: Experimental
    add_node_box(slide, E(6120000), E(4032000), E(1620000), E(648000),
                 ["③ 直接实验", "Apollo 15"],
                 NODE_ORANGE, Pt(9))
    add_label(slide, E(7920000), E(4032000), E(2700000), E(648000),
              "Pkg 6: 月球真空中\n锤子=羽毛 (质量比44:1)\n直接实验证实",
              Pt(8), RGBColor(0xBB, 0xBB, 0xBB))

    # Arrows converging
    add_arrow(slide, E(6930000), E(2736000), E(8640000), E(5040000))
    add_arrow(slide, E(6930000), E(3708000), E(8640000), E(5040000))
    add_arrow(slide, E(6930000), E(4680000), E(8640000), E(5040000))

    # Central conclusion node
    add_node_box(slide, E(7740000), E(5040000), E(3240000), E(900000),
                 ["5012: 真空等速", "belief = 0.95+", "三线汇聚 → 高确信度"],
                 NODE_GREEN, Pt(10))

    # Bottom note
    add_label(slide, RX, E(6048000), E(RW), E(360000),
              "旧理论不删除，belief 自然下降；全程只 add node+edge",
              Pt(9), RGBColor(0x99, 0x99, 0x99))


# ═══════════════════════════════════════════════════════════════
# Subtitle updates
# ═══════════════════════════════════════════════════════════════
SUBTITLES = {
    0: "伽利略绑球实验 — 概述",
    1: "Pkg 1: 亚里士多德的物理学 — 建立基线",
    2: "Pkg 2: 绑球悖论 — 纯推理推翻千年学说",
    3: "Pkg 3-4: 密度观察 → 真空等速预测",
    4: "Pkg 5: 牛顿力学 — 第二条独立矛盾线",
    5: "Pkg 6: Apollo 15 — 决定性实验",
}


# ═══════════════════════════════════════════════════════════════
# Main
# ═══════════════════════════════════════════════════════════════
def main():
    prs = Presentation("Gaia_Presentation.pptx")

    print(f"Total slides before: {len(prs.slides)}")

    # Step 1: Duplicate slide 10 (idx 10) twice to get 2 extra slides
    # They'll be appended at the end; we'll move them later
    duplicate_slide(prs, 10)  # will be new Pkg 5 slide
    duplicate_slide(prs, 10)  # will be new Pkg 6 slide

    total = len(prs.slides)
    print(f"Total slides after duplication: {total}")

    # Step 2: Move the two new slides (at end) to positions 13 and 14
    # Current: ... [9:overview] [10:pkg1-2] [11:pkg3-5] [12:pkg6] ... [new1] [new2]
    # Target:  ... [9:overview] [10:pkg1] [11:pkg2] [12:pkg3-4] [13:pkg5] [14:pkg6] ...
    # After dup, new slides are at indices total-2 and total-1

    # Move last slide (dup2) to index 13 (after current slide 12)
    move_slide(prs, total - 1, 13)
    # Move second-to-last (dup1, now at total-1 again since dup2 moved) to index 13
    move_slide(prs, total - 1, 13)

    print(f"Slides reordered. Total: {len(prs.slides)}")

    # Now Galileo slides are at indices 9-14 (6 slides)
    # The old slide at idx 10 had Pkg 1-2 content (will become Pkg 1)
    # The old slide at idx 11 had Pkg 3-5 content (will become Pkg 2)
    # The old slide at idx 12 had Pkg 6 content (will become Pkg 3-4)
    # New slides at idx 13, 14 are copies of idx 10 (will become Pkg 5, Pkg 6)

    # Step 3: Update each Galileo slide
    galileo_slides = [prs.slides[i] for i in range(9, 15)]

    # --- Slide 10 (idx 9): Overview ---
    slide_overview = galileo_slides[0]
    # Update the left text box (shape 3)
    tf = slide_overview.shapes[3].text_frame
    set_colored_text(tf, OVERVIEW_TEXT)
    # Update subtitle
    slide_overview.shapes[1].text_frame.paragraphs[0].runs[0].text = SUBTITLES[0]

    # --- Slide 11 (idx 10): Pkg 1 ---
    slide_pkg1 = galileo_slides[1]
    slide_pkg1.shapes[1].text_frame.paragraphs[0].runs[0].text = SUBTITLES[1]
    set_colored_text(slide_pkg1.shapes[3].text_frame, PKG1_CODE)
    build_diagram_pkg1(slide_pkg1)

    # --- Slide 12 (idx 11): Pkg 2 ---
    slide_pkg2 = galileo_slides[2]
    slide_pkg2.shapes[1].text_frame.paragraphs[0].runs[0].text = SUBTITLES[2]
    set_colored_text(slide_pkg2.shapes[3].text_frame, PKG2_CODE)
    build_diagram_pkg2(slide_pkg2)

    # --- Slide 13 (idx 12): Pkg 3-4 ---
    slide_pkg34 = galileo_slides[3]
    slide_pkg34.shapes[1].text_frame.paragraphs[0].runs[0].text = SUBTITLES[3]
    set_colored_text(slide_pkg34.shapes[3].text_frame, PKG34_CODE)
    build_diagram_pkg34(slide_pkg34)

    # --- Slide 14 (idx 13): Pkg 5 ---
    slide_pkg5 = galileo_slides[4]
    slide_pkg5.shapes[1].text_frame.paragraphs[0].runs[0].text = SUBTITLES[4]
    set_colored_text(slide_pkg5.shapes[3].text_frame, PKG5_CODE)
    build_diagram_pkg5(slide_pkg5)

    # --- Slide 15 (idx 14): Pkg 6 ---
    slide_pkg6 = galileo_slides[5]
    slide_pkg6.shapes[1].text_frame.paragraphs[0].runs[0].text = SUBTITLES[5]
    set_colored_text(slide_pkg6.shapes[3].text_frame, PKG6_CODE)
    build_diagram_pkg6(slide_pkg6)

    # Step 4: Save
    prs.save("Gaia_Presentation.pptx")
    print("Done! Galileo slides restructured (6 slides).")


if __name__ == "__main__":
    main()
