"""Add architecture overview slide as Slide 2 (after Title)."""

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from copy import deepcopy

NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"

# Colors
BOX_BLUE = RGBColor(0x1B, 0x5E, 0x9F)
BOX_TEAL = RGBColor(0x00, 0x69, 0x6B)
BOX_PURPLE = RGBColor(0x6A, 0x1B, 0x9A)
BOX_DARK = RGBColor(0x42, 0x42, 0x42)
ARROW_COLOR = RGBColor(0x90, 0xCA, 0xF9)
OUTPUT_COLOR = RGBColor(0x4E, 0x4E, 0x4E)
LABEL_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LABEL_LIGHT = RGBColor(0xB0, 0xB0, 0xB0)
LABEL_CYAN = RGBColor(0x5C, 0xC9, 0xF5)
LABEL_GOLD = RGBColor(0xFF, 0xD5, 0x4F)


def duplicate_slide(prs, source_idx):
    """Duplicate a slide to the end."""
    source = prs.slides[source_idx]
    new_slide = prs.slides.add_slide(source.slide_layout)
    sp_tree = new_slide.shapes._spTree
    for child in list(sp_tree)[2:]:
        sp_tree.remove(child)
    src_tree = source.shapes._spTree
    for child in list(src_tree)[2:]:
        sp_tree.append(deepcopy(child))
    src_bg = source._element.find(f"{{{NS_P}}}bg")
    if src_bg is not None:
        new_bg = new_slide._element.find(f"{{{NS_P}}}bg")
        if new_bg is not None:
            new_slide._element.remove(new_bg)
        new_slide._element.insert(0, deepcopy(src_bg))
    return new_slide


def move_slide(prs, from_idx, to_idx):
    """Move slide from from_idx to before to_idx."""
    sldIdLst = prs.slides._sldIdLst
    slides = list(sldIdLst)
    el = slides[from_idx]
    sldIdLst.remove(el)
    slides = list(sldIdLst)
    if to_idx >= len(slides):
        sldIdLst.append(el)
    else:
        slides[to_idx].addprevious(el)


def add_box(slide, left, top, width, height, lines, fill_color,
            font_size=Pt(11), font_color=LABEL_WHITE, first_bold=True):
    """Add a rounded rectangle with multi-line text."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = fill_color
    shape.line.width = Emu(12700)

    tf = shape.text_frame
    tf.word_wrap = True
    for i, line_text in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line_text
        p.font.size = font_size
        p.font.color.rgb = font_color
        p.font.bold = first_bold and (i == 0)
        p.space_after = Pt(2)
        p.space_before = Pt(2)
        p.alignment = PP_ALIGN.CENTER

    return shape


def add_label(slide, left, top, width, height, text,
              font_size=Pt(10), font_color=LABEL_LIGHT, bold=False):
    """Add a text label (no fill)."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.space_after = Pt(0)
    p.alignment = PP_ALIGN.CENTER
    return txBox


def add_arrow_line(slide, x1, y1, x2, y2, color=ARROW_COLOR):
    """Add a connector line (no arrowhead, just a line)."""
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)
    connector.line.color.rgb = color
    connector.line.width = Emu(25400)  # 2pt
    return connector


def add_output_pill(slide, left, top, width, height, text, fill_color):
    """Add a small rounded rectangle for output types."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = fill_color
    shape.line.width = Emu(6350)

    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = LABEL_WHITE
    p.font.bold = False
    p.space_after = Pt(0)
    p.alignment = PP_ALIGN.CENTER
    return shape


def clear_all_shapes(slide):
    """Remove all shapes from a slide."""
    sp_tree = slide.shapes._spTree
    for child in list(sp_tree)[2:]:
        sp_tree.remove(child)


def main():
    prs = Presentation("Gaia_Presentation.pptx")
    total = len(prs.slides)
    print(f"Starting with {total} slides")

    # Duplicate slide 2 (idx 1, Part 1 layout) as template for dark background
    duplicate_slide(prs, 1)
    # Move to index 1 (after title)
    move_slide(prs, len(prs.slides) - 1, 1)
    print(f"Added slide. Now: {len(prs.slides)} slides")

    # Get reference to the new slide at index 1
    slide = prs.slides[1]

    # Clear all shapes from template
    clear_all_shapes(slide)

    # ═══════════════════════════════════════════════════════════════
    # Layout constants (slide: 12192000 x 6858000 EMU)
    # ═══════════════════════════════════════════════════════════════

    # Title area
    add_label(slide, 720000, 180000, 3600000, 360000,
              "Overview", Pt(13), LABEL_LIGHT, bold=False)
    add_label(slide, 720000, 420000, 10440000, 600000,
              "Gaia 全景架构", Pt(28), LABEL_WHITE, bold=True)

    # Separator line
    slide.shapes.add_connector(1,
        Emu(720000), Emu(1100000),
        Emu(11160000), Emu(1100000)
    ).line.color.rgb = RGBColor(0x55, 0x55, 0x55)

    # ═══════════════════════════════════════════════════════════════
    # Main architecture boxes (3 boxes, horizontal)
    # ═══════════════════════════════════════════════════════════════
    BOX_W = 2800000
    BOX_H = 1100000
    BOX_Y = 1500000
    GAP = 900000

    # Calculate x positions (centered on slide)
    total_w = 3 * BOX_W + 2 * GAP
    start_x = (12192000 - total_w) // 2

    x1 = start_x
    x2 = start_x + BOX_W + GAP
    x3 = start_x + 2 * (BOX_W + GAP)

    # Box 1: Research Agent
    add_box(slide, x1, BOX_Y, BOX_W, BOX_H,
            ["Research Agent", "(AI 科研助手)"],
            BOX_PURPLE, Pt(14))

    # Box 2: Gaia CLI
    add_box(slide, x2, BOX_Y, BOX_W, BOX_H,
            ["Gaia CLI", "(local)"],
            BOX_TEAL, Pt(14))

    # Box 3: Gaia Server
    add_box(slide, x3, BOX_Y, BOX_W, BOX_H,
            ["Gaia Server", "(Large Knowledge Model)"],
            BOX_BLUE, Pt(14))

    # ═══════════════════════════════════════════════════════════════
    # Arrows between boxes (bidirectional, using two lines with gap)
    # ═══════════════════════════════════════════════════════════════
    arrow_y = BOX_Y + BOX_H // 2

    # Arrow 1: Agent ←→ CLI
    add_arrow_line(slide, x1 + BOX_W, arrow_y - 20000,
                   x2, arrow_y - 20000, ARROW_COLOR)
    add_arrow_line(slide, x2, arrow_y + 20000,
                   x1 + BOX_W, arrow_y + 20000, ARROW_COLOR)

    # Arrow 2: CLI ←→ Server
    add_arrow_line(slide, x2 + BOX_W, arrow_y - 20000,
                   x3, arrow_y - 20000, ARROW_COLOR)
    add_arrow_line(slide, x3, arrow_y + 20000,
                   x2 + BOX_W, arrow_y + 20000, ARROW_COLOR)

    # ═══════════════════════════════════════════════════════════════
    # Annotation labels between boxes
    # ═══════════════════════════════════════════════════════════════
    label_w = GAP + 200000
    label_y = BOX_Y + BOX_H + 60000

    # Between Agent and CLI
    add_label(slide, x1 + BOX_W - 100000, label_y - 80000, label_w, 400000,
              "external memory\nfor agents",
              Pt(9), LABEL_CYAN, bold=False)

    # Between CLI and Server
    add_label(slide, x2 + BOX_W - 100000, label_y - 80000, label_w, 400000,
              "commit → review\n→ merge",
              Pt(9), LABEL_CYAN, bold=False)

    # ═══════════════════════════════════════════════════════════════
    # Description labels below each box
    # ═══════════════════════════════════════════════════════════════
    desc_y = BOX_Y + BOX_H + 420000

    add_label(slide, x1 - 100000, desc_y, BOX_W + 200000, 500000,
              "Agent 读写知识\n作为 external memory",
              Pt(9), LABEL_LIGHT, bold=False)

    add_label(slide, x2 - 100000, desc_y, BOX_W + 200000, 500000,
              "knowledge packages\ninit / commit / publish",
              Pt(9), LABEL_LIGHT, bold=False)

    add_label(slide, x3 - 100000, desc_y, BOX_W + 200000, 500000,
              "十亿级推理超图\nBP · Contradiction",
              Pt(9), LABEL_LIGHT, bold=False)

    # ═══════════════════════════════════════════════════════════════
    # Output section: arrow down from Server → output types
    # ═══════════════════════════════════════════════════════════════
    server_center_x = x3 + BOX_W // 2
    arrow_top = desc_y + 400000
    arrow_bottom = arrow_top + 500000

    # Vertical arrow from server area down
    add_arrow_line(slide, server_center_x, arrow_top,
                   server_center_x, arrow_bottom, ARROW_COLOR)

    # "知识输出" label
    add_label(slide, server_center_x + 40000, arrow_top + 100000, 1200000, 300000,
              "知识输出 ↓", Pt(10), LABEL_GOLD, bold=True)

    # Output pills (2 rows of 3)
    outputs = ["论文", "百科", "综述", "科普", "教材", "FAQ"]
    pill_w = 1000000
    pill_h = 340000
    pill_gap = 100000
    row1_count = 3
    row1_w = row1_count * pill_w + (row1_count - 1) * pill_gap
    pill_start_x = server_center_x - row1_w // 2
    pill_y1 = arrow_bottom + 100000
    pill_y2 = pill_y1 + pill_h + pill_gap

    output_colors = [
        RGBColor(0x1B, 0x5E, 0x9F),  # blue
        RGBColor(0x2E, 0x7D, 0x32),  # green
        RGBColor(0x6A, 0x1B, 0x9A),  # purple
        RGBColor(0xE6, 0x51, 0x00),  # orange
        RGBColor(0x00, 0x69, 0x6B),  # teal
        RGBColor(0x55, 0x55, 0x55),  # gray
    ]

    for i, (label, color) in enumerate(zip(outputs, output_colors)):
        row = i // 3
        col = i % 3
        px = pill_start_x + col * (pill_w + pill_gap)
        py = pill_y1 if row == 0 else pill_y2
        add_output_pill(slide, px, py, pill_w, pill_h, label, color)

    # Save
    prs.save("Gaia_Presentation.pptx")
    print(f"Saved. Final: {len(prs.slides)} slides")

    # Verify
    prs2 = Presentation("Gaia_Presentation.pptx")
    for i in range(min(5, len(prs2.slides))):
        slide = prs2.slides[i]
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text[:50]
                if t.strip():
                    texts.append(t)
        first = texts[0] if texts else "(empty)"
        print(f"  {i+1}. {first}")


if __name__ == "__main__":
    main()
