"""Add syntax highlighting to code blocks in Galileo + Einstein slides."""

from pptx import Presentation
from copy import deepcopy
from lxml import etree

NSMAP = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}

# Color palette (hex strings, for dark backgrounds)
WHITE = "E0E0E0"
CYAN = "5CC9F5"  # $ gaia commands
GRAY = "808888"  # comments
GOLD = "FFD54F"  # ★ key conclusions
RED = "FF8A80"  # belief ↓
GREEN = "69F0AE"  # belief ↑
ORANGE = "FFAB40"  # contradiction
HEADER = "FFFFFF"  # package headers


def colorize_line(line):
    """Return list of (text, color_hex, bold) segments for a line."""
    stripped = line.strip()

    if not stripped:
        return [("", WHITE, False)]

    # Package headers: # ═══ Pkg N: ... ═══
    if "═══" in stripped:
        return [(line, HEADER, True)]

    # Comments (but not package headers)
    if stripped.startswith("#") and "═══" not in stripped:
        # Check for special keywords in comments
        if "CONTRADICTION" in stripped or "矛盾" in stripped:
            return [(line, ORANGE, False)]
        return [(line, GRAY, False)]

    # CLI commands: $ gaia ...
    if stripped.startswith("$ gaia"):
        parts = []
        # Split command from trailing comment
        if "  #" in line:
            idx = line.index("  #")
            parts.append((line[:idx], CYAN, False))
            parts.append((line[idx:], GRAY, False))
        elif "# " in line and not line.strip().startswith("#"):
            idx = line.index("# ")
            parts.append((line[:idx], CYAN, False))
            parts.append((line[idx:], GRAY, False))
        else:
            parts.append((line, CYAN, False))
        return parts

    # Key conclusions with ★
    if "★" in line:
        return [(line, GOLD, True)]

    # Contradiction lines
    if "contradiction" in stripped.lower() or "CONTRADICTION" in stripped:
        return [(line, ORANGE, True)]

    # Output arrows: → created node ...
    if stripped.startswith("→") or stripped.startswith("→"):
        if "★" in line:
            return [(line, GOLD, True)]
        return [(line, GREEN, False)]

    # Belief results: lines with ↓ and/or ↑
    has_down = "↓" in line
    has_up = "↑" in line
    if has_down and has_up:
        # Mixed line — split into segments by whitespace groups
        segments = []
        parts = line.split("  ")  # Split by double-space
        for part in parts:
            part_stripped = part.strip()
            if not part_stripped:
                continue
            if "↓" in part:
                segments.append(("  " + part_stripped, RED, False))
            elif "↑" in part:
                segments.append(("  " + part_stripped, GREEN, False))
            else:
                segments.append(("  " + part_stripped, WHITE, False))
        return segments if segments else [(line, WHITE, False)]
    if has_down:
        return [(line, RED, False)]
    if has_up:
        return [(line, GREEN, False)]

    return [(line, WHITE, False)]


def make_run_xml(text, color_hex, bold=False):
    """Create an <a:r> element with colored text."""
    ns = NSMAP["a"]
    r = etree.Element(f"{{{ns}}}r")
    rPr = etree.SubElement(r, f"{{{ns}}}rPr")
    rPr.set("lang", "en-US")
    rPr.set("dirty", "0")
    if bold:
        rPr.set("b", "1")

    solidFill = etree.SubElement(rPr, f"{{{ns}}}solidFill")
    srgbClr = etree.SubElement(solidFill, f"{{{ns}}}srgbClr")
    srgbClr.set("val", color_hex)

    t = etree.SubElement(r, f"{{{ns}}}t")
    t.text = text
    # Preserve leading/trailing spaces
    if text and (text[0] == " " or text[-1] == " "):
        t.set("{http://www.w3.org/XML/1998/namespace}space", "preserve")

    return r


def set_colored_text(text_frame, lines):
    """Replace text frame content with syntax-highlighted lines."""
    # Get reference paragraph for spacing/format
    ref_p = text_frame.paragraphs[0]._p
    ref_pPr = ref_p.find("a:pPr", NSMAP)
    ref_pPr_xml = deepcopy(ref_pPr) if ref_pPr is not None else None

    # Clear all paragraphs
    txBody = text_frame._txBody
    for p in txBody.findall(f"{{{NSMAP['a']}}}p"):
        txBody.remove(p)

    ns = NSMAP["a"]

    for i, line in enumerate(lines):
        p = etree.SubElement(txBody, f"{{{ns}}}p")

        # Add paragraph properties (spacing etc.) from reference
        if ref_pPr_xml is not None:
            pPr = deepcopy(ref_pPr_xml)
            # Only first line gets center alignment
            if i > 0 and "algn" in pPr.attrib:
                del pPr.attrib["algn"]
            p.insert(0, pPr)

        # Colorize and add runs
        segments = colorize_line(line)
        for text, color, bold in segments:
            r = make_run_xml(text, color, bold)
            p.append(r)


def main():
    prs = Presentation("Gaia_Presentation.pptx")

    # ================================================================
    # Slide 11 (index 10): Galileo Pkg 1 + Pkg 2
    # ================================================================
    slide11 = prs.slides[10]
    set_colored_text(
        slide11.shapes[3].text_frame,
        [
            "$ gaia init galileo_tied_balls",
            "",
            "# ═══ Pkg 1: aristotle_physics ═══",
            '$ gaia node add "越重的物体下落越快" \\',
            "    --prior 0.70 --type paper-extract   → 5001",
            '$ gaia node add "石头比树叶落得快" \\',
            "    --prior 0.90                        → 5002",
            "$ gaia edge add --tail 5001,5002 \\",
            "    --head 5003 --type abstraction",
            '  → 5003 "v ∝ W" prior=0.70',
            '$ gaia commit -m "aristotle_physics"',
            "$ gaia propagate",
            "  5003 (v ∝ W): belief = 0.70",
            "",
            "# ═══ Pkg 2: galileo1638_tied_balls ═══",
            '$ gaia node add "重球H绑轻球L"        → 5004',
            "$ gaia edge add --tail 5003,5004 --head 5005 \\",
            '    --type deduction  → "推导A: HL更慢"',
            "$ gaia edge add --tail 5003,5004 --head 5006 \\",
            '    --type deduction  → "推导B: HL更快"',
            "$ gaia edge add --tail 5005,5006 \\",
            "    --type contradiction",
            "  # 同一物体不可能既更快又更慢 → 矛盾!",
            "$ gaia commit && gaia propagate",
            "  5003 (v∝W): 0.70→0.35 ↓  矛盾回传",
            "  5008 (v∝W错误): belief=0.82 ↑",
        ],
    )

    # ================================================================
    # Slide 12 (index 11): Galileo Pkg 3 + Pkg 4 + Pkg 5
    # ================================================================
    slide12 = prs.slides[11]
    set_colored_text(
        slide12.shapes[3].text_frame,
        [
            "# ═══ Pkg 3: medium_density ═══",
            '$ gaia node add "水中落速差异更大" --prior 0.90',
            '$ gaia node add "介质密度↓ 差异↓" --prior 0.85',
            "$ gaia edge add --tail 5009,5010 --head 5011 \\",
            "    --type deduction",
            '  → 5011 "★ 结论①: 空气阻力才是混淆因素"',
            "$ gaia commit && gaia propagate",
            "  5003: 0.35→0.28 ↓  替代解释削弱旧定律",
            "",
            "# ═══ Pkg 4: vacuum_prediction ═══",
            "$ gaia edge add --tail 5008,5011 --head 5012",
            '  → 5012 "★ 结论②: 真空中所有物体等速下落"',
            '$ gaia node add "斜面实验: 不同球同时到底" \\',
            "    --prior 0.90",
            "$ gaia commit && gaia propagate",
            "  5012 (真空等速): 0.00→0.78 ↑",
            "",
            "# ═══ Pkg 5: newton_principia ═══",
            '$ gaia node add "F=ma" "F=mg" --prior 0.95',
            "$ gaia edge add --tail 5015,5016 --head 5017",
            '  → 5017 "a=g 质量约掉了!"',
            "$ gaia edge add --tail 5003,5017 \\",
            "    --type contradiction  # 第二条矛盾线!",
            "$ gaia commit && gaia propagate",
            "  5003→0.12↓↓  5017→0.93↑  5012→0.87↑",
        ],
    )

    # ================================================================
    # Slide 13 (index 12): Galileo Pkg 6 + Summary
    # ================================================================
    slide13 = prs.slides[12]
    set_colored_text(
        slide13.shapes[3].text_frame,
        [
            "# ═══ Pkg 6: apollo15_feather_drop ═══",
            '$ gaia node add "月球: ≈3e-12 atm" \\',
            "    --prior 0.99                      → 5018",
            '$ gaia node add "锤子=羽毛 同时落地" \\',
            "    --prior 0.99  # 质量比 44:1!      → 5019",
            "$ gaia edge add --tail 5018,5019 \\",
            "    --head 5020 --type deduction",
            "",
            '$ gaia commit -m "apollo15"',
            "$ gaia propagate",
            "  5003 (v∝W): 0.12→0.05 ↓  几乎归零",
            "  5012 (真空等速): 0.87→0.95 ↑  三线汇聚",
            "  5017 (a=g): 0.93→0.96 ↑  理论+实验",
            "  5020 (月球): belief=0.98 ↑  决定性",
            "",
            "$ gaia test     # 验证 beliefs.yaml",
            "$ gaia publish  # → 远程 registry",
        ],
    )

    # ================================================================
    # Slide 15 (index 14): Einstein Pkg 1-3
    # ================================================================
    slide15 = prs.slides[14]
    set_colored_text(
        slide15.shapes[3].text_frame,
        [
            "$ gaia init einstein_elevator",
            "",
            "# ═══ Pkg 1: prior_knowledge ═══",
            '$ gaia node add "F=mᵢa" --prior 0.95        → 6001',
            '$ gaia node add "F=GMm/r²" --prior 0.95     → 6002',
            '$ gaia node add "光的微粒说" --prior 0.50     → 6003',
            '$ gaia node add "Maxwell EM" --prior 0.90    → 6005',
            '$ gaia node add "Eötvös mᵢ=mᵍ" --prior 0.95 → 6006',
            "$ gaia edge add --tail 6001,6002,6003 \\",
            "    --head 6004 --type deduction  P=0.60",
            '  → 6004 "Soldner: 0.87″" prior=0.60',
            "",
            "# ═══ Pkg 2: equivalence_principle ═══",
            '$ gaia node add "电梯思想实验" --prior 1.0  → 6007',
            "$ gaia edge add --tail 6006,6007 --head 6008",
            '  → 6008 "等效原理" belief≈0.85',
            "$ gaia edge add --tail 6008,6005 --head 6009",
            '  → 6009 "光在引力场弯曲"',
            "",
            "# ═══ Pkg 3: 1911_light_deflection ═══",
            "$ gaia edge add --tail 6009 --head 6010",
            '  → 6010 "Einstein: 0.87″" belief≈0.80',
            "$ gaia commit && gaia propagate",
            "# ⚠ 6004 = 6010 = 0.87″  此时无矛盾!",
        ],
    )

    # ================================================================
    # Slide 16 (index 15): Einstein Pkg 4
    # ================================================================
    slide16 = prs.slides[15]
    set_colored_text(
        slide16.shapes[3].text_frame,
        [
            "# ═══ Pkg 4: general_relativity ═══",
            '$ gaia node add "GR: Gμν=8πGTμν" \\',
            "    --prior 0.85                      → 6012",
            "$ gaia edge add --tail 6012,6008 --head 6013",
            '  → 6013 "光沿零测地线 (时间+空间弯曲)"',
            "$ gaia edge add --tail 6013 --head 6014",
            '  → 6014 "GR: 1.75″ = 2×0.87″"',
            '$ gaia node add "水星进动 43″/世纪" \\',
            "    --prior 0.90                      → 6015",
            "$ gaia edge add --tail 6012 --head 6015",
            "",
            "# 关键: GR vs Newton 定量矛盾!",
            "$ gaia edge add --tail 6014,6004 \\",
            "    --type contradiction",
            "  # 1.75″ vs 0.87″ — 精确 2× 差异",
            "",
            '$ gaia commit -m "general_relativity"',
            "$ gaia propagate",
            "  6014 (GR 1.75″): belief=0.85",
            "  6004 (Newton 0.87″): 0.60→0.40 ↓",
        ],
    )

    # ================================================================
    # Slide 17 (index 16): Einstein Pkg 5
    # ================================================================
    slide17 = prs.slides[16]
    set_colored_text(
        slide17.shapes[3].text_frame,
        [
            "# ═══ Pkg 5: eddington1919_eclipse ═══",
            '$ gaia node add "日食远征(Sobral+Príncipe)" \\',
            "    --prior 0.95                      → 6017",
            '$ gaia node add "观测 1.61±0.30″ / 1.98±0.16″" \\',
            "    --prior 0.90                      → 6018",
            "$ gaia edge add --tail 6017,6018 --head 6019",
            '  → 6019 "观测支持 GR, 排除 Newton"',
            "",
            "# 观测 vs 牛顿: 实验排除",
            "$ gaia edge add --tail 6019,6004 \\",
            "    --type contradiction  # ~1.7″ vs 0.87″",
            "",
            '$ gaia commit -m "eddington1919"',
            "$ gaia publish  # → 远程 registry",
            "$ gaia review <commit_id>",
            "$ gaia merge <commit_id>",
            "  6004: 0.40→0.10 ↓↓  6014: →0.95 ↑↑",
            "  6002 (Newton引力): 0.95→0.80 ↓ 降级",
        ],
    )

    prs.save("Gaia_Presentation.pptx")
    print("Colorized all code blocks")


if __name__ == "__main__":
    main()
