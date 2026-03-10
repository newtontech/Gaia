"""Update Gaia presentation slides to match CLI design doc."""

from pptx import Presentation
from copy import deepcopy
from lxml import etree


def clear_and_set_text(text_frame, lines):
    """Replace all text in a text_frame with new lines, preserving first paragraph's formatting."""
    # Save reference paragraph XML for formatting
    ref_para = text_frame.paragraphs[0]._p
    ref_xml = deepcopy(ref_para)

    # Remove all existing paragraphs
    txBody = text_frame._txBody
    for p in txBody.findall(
        "{http://schemas.openxmlformats.org/drawingml/2006/main}p"
    ):
        txBody.remove(p)

    # Add new paragraphs
    nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    for i, line in enumerate(lines):
        new_p = deepcopy(ref_xml)
        # Clear runs from the template
        for r in new_p.findall("a:r", nsmap):
            new_p.remove(r)
        # Remove alignment for non-header lines
        pPr = new_p.find("a:pPr", nsmap)
        if pPr is not None and i > 0:
            if "algn" in pPr.attrib:
                del pPr.attrib["algn"]

        # Create a new run with the text
        r_elem = deepcopy(ref_xml.findall("a:r", nsmap)[0]) if ref_xml.findall("a:r", nsmap) else None
        if r_elem is not None:
            r_elem.find("a:t", nsmap).text = line
        else:
            r_elem = etree.SubElement(new_p, f'{{{nsmap["a"]}}}r')
            t_elem = etree.SubElement(r_elem, f'{{{nsmap["a"]}}}t')
            t_elem.text = line

        new_p.append(r_elem)
        txBody.append(new_p)


def update_table_cell(table, row, col, text):
    """Update text in a table cell."""
    cell = table.cell(row, col)
    for para in cell.text_frame.paragraphs:
        for run in para.runs:
            run.text = ""
    cell.text_frame.paragraphs[0].runs[0].text = text if cell.text_frame.paragraphs[0].runs else text
    # Fallback: set via text property
    if not cell.text_frame.paragraphs[0].runs:
        cell.text = text
    else:
        cell.text_frame.paragraphs[0].runs[0].text = text


def main():
    prs = Presentation("Gaia_Presentation.pptx")

    # ================================================================
    # Slide 10 (index 9): Galileo overview — update belief trajectory
    # ================================================================
    slide10 = prs.slides[9]
    for shape in slide10.shapes:
        if shape.has_text_frame:
            full_text = shape.text_frame.text
            # Update the belief trajectory values
            if "0.15 ↓" in full_text:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.text = run.text.replace("0.15 ↓", "0.28 ↓")
            if "0.08 ↓" in full_text:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.text = run.text.replace("0.08 ↓", "0.12 ↓")

    # ================================================================
    # Slide 11 (index 10): Galileo Pkg 1-2 — code block
    # ================================================================
    slide11 = prs.slides[10]
    code_block = slide11.shapes[3]  # Rounded Rectangle 4
    clear_and_set_text(code_block.text_frame, [
        "# packages/aristotle_physics.yaml",
        'node_5001 "自然运动学说"     prior=0.70',
        'node_5002 "石头快于叶子"     prior=0.90',
        'node_5003 "v ∝ W 定律"      prior=0.70',
        "edge 5001: [5001,5002] → [5003]  abstraction",
        "",
        "# packages/galileo1638_tied_balls.yaml",
        'node_5004 "绑球设定 H+L"    prior=0.99',
        'node_5005 "推导A: 组合更慢"   prior=0.90',
        'node_5006 "推导B: 组合更快"   prior=0.90',
        "edge 5002: [5003,5004] → [5005]  deduction",
        "edge 5003: [5003,5004] → [5006]  deduction",
        "edge 5004 [CONTRADICTION]: [5005] ↔ [5006]",
        "",
        "$ gaia build && gaia commit -m \"Pkg 1-2\"",
        "$ gaia propagate",
        "  5003 (v ∝ W): 0.70 → 0.35 ↓",
        "  5008 (定律错误): belief = 0.82 ↑",
    ])

    # ================================================================
    # Slide 12 (index 11): Galileo Pkg 3-5 — code block
    # ================================================================
    slide12 = prs.slides[11]
    code_block = slide12.shapes[3]  # Rounded Rectangle 4
    clear_and_set_text(code_block.text_frame, [
        "# packages/galileo1638_medium_density.yaml",
        'node_5009 "密度递减观察"     prior=0.90',
        'node_5010 "速差随密度↓"     prior=0.85',
        'node_5011 "空气阻力假说"     prior=0.80',
        "edge 5006: [5009,5010] → [5011]  deduction",
        "",
        "# packages/galileo1638_vacuum_prediction.yaml",
        'node_5012 "真空等速预测"     prior=0.85',
        "edge 5007: [5008,5011] → [5012]  deduction",
        "",
        "# packages/newton1687_principia.yaml",
        'node_5015 "F = ma"          prior=0.95',
        'node_5016 "F = mg"          prior=0.95',
        'node_5017 "∴ a = g"         prior=0.95',
        "edge 5010: [5015,5016] → [5017]  deduction",
        "edge 5011 [CONTRADICTION]: [5003] ↔ [5017]",
        "",
        "$ gaia commit -m \"Pkg 3-5\" && gaia propagate",
        "  5003→0.12↓  5012→0.87↑  5017→0.93↑",
    ])

    # Update diagram labels on slide 12
    for shape in slide12.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text
            if "0.35 → 0.08" in text:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.text = run.text.replace("0.35 → 0.08", "0.35 → 0.12")
            if "0.85 → 0.88" in text:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.text = run.text.replace("0.85 → 0.88", "0.78 → 0.87")
            if "从 0.35 降至 0.08" in text:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.text = run.text.replace("从 0.35 降至 0.08", "从 0.35 降至 0.12")

    # ================================================================
    # Slide 13 (index 12): Galileo Pkg 6 — code block + table
    # ================================================================
    slide13 = prs.slides[12]
    code_block = slide13.shapes[3]  # Rounded Rectangle 4
    clear_and_set_text(code_block.text_frame, [
        "# packages/apollo15_1971_feather_drop.yaml",
        'node_5018 "月球真空环境"     prior=0.99',
        'node_5019 "锤=羽同时落"     prior=0.99',
        'node_5020 "等速落体确认"     prior=0.99',
        "edge 5012: [5018,5019] → [5020]  P=0.99",
        "edge 5013: [5020] → [5012]      P=0.99",
        "",
        "$ gaia commit -m \"Pkg 6\"",
        "$ gaia propagate",
        "$ gaia publish  # → 推送到远程 registry",
    ])

    # Update table values on slide 13
    table_shape = slide13.shapes[4]  # Table 5
    if table_shape.has_table:
        table = table_shape.table
        # Row 3 (Pkg 3-4): col 1 = v∝W value, col 2 = a=g value
        update_table_cell(table, 3, 1, "0.28 ↓")
        update_table_cell(table, 3, 2, "0.78")
        # Row 4 (Pkg 5 Newton): col 1 = v∝W value
        update_table_cell(table, 4, 1, "0.12 ↓")

    # ================================================================
    # Slide 15 (index 14): Einstein Pkg 1-3 — code block
    # ================================================================
    slide15 = prs.slides[14]
    code_block = slide15.shapes[3]  # Rounded Rectangle 4
    clear_and_set_text(code_block.text_frame, [
        "# packages/prior_knowledge.yaml",
        'node_6001 "F = mᵢ·a"        prior=0.95',
        'node_6002 "F = GMm/r²"       prior=0.95',
        'node_6003 "光的微粒说"         prior=0.50',
        'node_6004 "Soldner: 0.87″"   prior=0.60',
        'node_6005 "Maxwell EM"        prior=0.90',
        'node_6006 "Eötvös mᵢ=mᵍ"    prior=0.95',
        "edge 6001: [6001,6002,6003] → [6004]  P=0.60",
        "",
        "# packages/einstein1907_equivalence.yaml",
        'node_6007 "电梯思想实验"       prior=1.00',
        'node_6008 "等效原理"           prior=0.85',
        'node_6009 "光在引力场弯曲"     prior=0.85',
        "edge 6002: [6006,6007] → [6008]  P=0.85",
        "edge 6003: [6008,6005] → [6009]  P=0.85",
        "",
        "# packages/einstein1911_light_deflection.yaml",
        'node_6010 "Einstein: 0.87″"  prior=0.80',
        "edge 6004: [6009] → [6010]    P=0.80",
        "",
        "$ gaia commit -m \"Pkg 1-3\" && gaia propagate",
        "# ⚠ 6004 = 6010 = 0.87″  此时无矛盾!",
    ])

    # ================================================================
    # Slide 16 (index 15): Einstein Pkg 4 — code block
    # ================================================================
    slide16 = prs.slides[15]
    code_block = slide16.shapes[3]  # Rounded Rectangle 4
    clear_and_set_text(code_block.text_frame, [
        "# packages/einstein1915_general_relativity.yaml",
        'node_6012 "GR: Gμν=8πGTμν"   prior=0.85',
        'node_6013 "光沿零测地线"        prior=0.85',
        "  (时间弯曲 + 空间弯曲 = 2个分量)",
        'node_6014 "GR预测: 1.75″"     prior=0.85',
        "  (= 2 × 0.87″, 空间弯曲贡献一半)",
        'node_6015 "水星进动 43″/世纪"   prior=0.90',
        "",
        "edge 6006: [6012,6008] → [6013]  P=0.85",
        "edge 6007: [6013] → [6014]       P=0.85",
        "edge 6008 [CONTRADICTION]:",
        "  [6014: 1.75″] ↔ [6004: 0.87″]  # 精确 2×!",
        "edge 6009: [6012] → [6015]       P=0.90",
        "",
        "$ gaia commit -m \"Pkg 4\" && gaia propagate",
        "  6014 (GR 1.75″): belief = 0.85",
        "  6004 (Newton 0.87″): 0.60 → 0.40 ↓",
    ])

    # ================================================================
    # Slide 17 (index 16): Einstein Pkg 5 — code block
    # ================================================================
    slide17 = prs.slides[16]
    code_block = slide17.shapes[3]  # Rounded Rectangle 4
    clear_and_set_text(code_block.text_frame, [
        "# packages/eddington1919_solar_eclipse.yaml",
        'node_6017 "日食远征(两队)"     prior=0.95',
        'node_6018 "观测 1.61±0.30″ /  prior=0.90',
        '          1.98±0.16″"',
        'node_6019 "支持GR排除Newton"   prior=0.90',
        'node_6020 "GR确认为优越理论"    prior=0.90',
        'node_6021 "Newton降级为近似"    prior=0.90',
        "",
        "edge 6011: [6017,6018] → [6019]  P=0.90",
        "edge 6014 [CONTRADICTION]:",
        "  [6019: ~1.7″] ↔ [6004: 0.87″]",
        "",
        "$ gaia publish  # 提交到远程 registry",
        "$ gaia review <commit_id>  # LLM 审查",
        "$ gaia merge <commit_id>   # 合并 + BP",
    ])

    # ================================================================
    # Slide 18 (index 17): Knowledge Package — add CLI context
    # ================================================================
    slide18 = prs.slides[17]
    # Update the text box with bullet points (Shape 3)
    text_shape = slide18.shapes[3]  # TextBox 4
    clear_and_set_text(text_shape.text_frame, [
        "类比 Cargo / npm：论文 = crate / package",
        "gaia.toml 声明包元数据 + 依赖",
        "packages/*.yaml 定义节点和边",
        "可版本化、可声明依赖、可追踪变更",
        "",
        "  [package]",
        '  name = "galileo_tied_balls"',
        '  version = "1.0.0"',
        "",
        "  [[packages]]",
        '  name = "aristotle_physics"',
        "  order = 1",
        "",
        "  [[packages]]",
        '  name = "galileo1638_tied_balls"',
        "  order = 2",
        '  depends_on = ["aristotle_physics"]',
    ])

    # ================================================================
    # Slide 20 (index 19): Replace "Environment" with CLI Workflow
    # ================================================================
    slide20 = prs.slides[19]

    # Update title
    title_shape = slide20.shapes[1]  # TextBox 2
    clear_and_set_text(title_shape.text_frame, [
        "Gaia CLI 工作流",
    ])

    # Update body text
    body_shape = slide20.shapes[3]  # TextBox 4
    clear_and_set_text(body_shape.text_frame, [
        "声明式包文件 + 交互式命令，双模混合",
        "本地优先（LanceDB + Kuzu），远程可选",
        "远程发布完整暴露三步流程:",
        "  publish → review → merge",
        "",
        "$ gaia init galileo_tied_balls",
        "$ gaia build       # 校验包结构",
        "$ gaia commit      # 提交到本地图",
        "$ gaia propagate   # 本地 BP",
        "$ gaia test        # belief 断言",
        "$ gaia publish     # → 远程 registry",
        "$ gaia review <id> # LLM 审查",
        "$ gaia merge <id>  # 合并 + 远程 BP",
    ])

    # Update branch labels
    for shape in slide20.shapes:
        if shape.has_text_frame:
            if shape.text_frame.text.strip() == "main":
                clear_and_set_text(shape.text_frame, ["local graph"])
            elif shape.text_frame.text.strip() == "experiment branch":
                clear_and_set_text(shape.text_frame, ["publish → review → merge → remote registry"])

    # ================================================================
    # Save
    # ================================================================
    output = "Gaia_Presentation.pptx"
    prs.save(output)
    print(f"Saved updated presentation to {output}")


if __name__ == "__main__":
    main()
